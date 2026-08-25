"""
Assemble a workshop slide deck (16:9) from the same figures and numbers used in
generate_report.py. This is a *deck builder*, not a modeling script: it reads
figure files and reproduces the report's own prose as condensed bullets: it
does not recompute anything.

Every slide also carries a PowerPoint speaker note (presenter view / View >
Notes Page) with talking points, transitions, and — where relevant — which
notebook to have open. Section dividers additionally show which notebook(s)
from notebooks/00-07 run during that section, so the deck doubles as a
run sheet for the live workshop.

Run:  python scripts/generate_slides.py
Writes: report/Beyond_Interoperability_Slides.pptx
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.join(os.path.dirname(__file__), "..")
FIG = os.path.join(ROOT, "figures")
OUT_DIR = os.path.join(ROOT, "report")
OUT = os.path.join(OUT_DIR, "Beyond_Interoperability_Slides.pptx")
os.makedirs(OUT_DIR, exist_ok=True)

# ----------------------------------------------------------------------------
# Palette (matches generate_report.py) and layout constants
# ----------------------------------------------------------------------------

INK = RGBColor(0x1a, 0x1a, 0x1a)
INK_2 = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x2a, 0x78, 0xd6)
ACCENT_DARK = RGBColor(0x1c, 0x5c, 0xab)
MUTED = RGBColor(0x60, 0x60, 0x60)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT = RGBColor(0xf0, 0xf4, 0xfa)
PLACEHOLDER = RGBColor(0xc0, 0x3a, 0x2b)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

DECK_TITLE = "Beyond Interoperability"
_page_counter = {"n": 0}

# ----------------------------------------------------------------------------
# The workshop notebook sequence — run in this numeric order during the
# hands-on portions of the workshop. Kept as data so both the summary slide
# and the section-divider badges stay in sync with one source of truth.
# ----------------------------------------------------------------------------

NOTEBOOKS = [
    ("00", "input_data_exploration.ipynb",
     "Load the 4 raw hospital files, inspect columns, first-look plots",
     "Section 5 — Our Dataset"),
    ("01", "site_heterogeneity.ipynb",
     "Missingness, label shift, covariate shift, domain-classifier alarm",
     "Section 5 — Our Dataset"),
    ("02", "conformal_basics.ipynb",
     "Build a LAC conformal predictor from scratch in ~5 lines of NumPy",
     "Section 7 / 8 — Conformal Prediction & the Math"),
    ("03", "federated_training.ipynb",
     "Implement FedAvg from scratch; train the shared model across 3 sites",
     "Section 9 — Pipeline (Training)"),
    ("04", "conformal_under_site_shift.ipynb",
     "Calibrate on Cleveland only, deploy everywhere — the coverage break",
     "Section 7 & 9 — Headline Result / Deployment"),
    ("05", "recreating_paper_figures.ipynb",
     "Recreate the Angelopoulos & Bates illustrations on this project's data",
     "Section 8 — The Mathematics"),
    ("06", "chest_pain_multiclass.ipynb",
     "Repeat the full pipeline on a second, independent target",
     "Section 9 — Robustness (chest-pain task)"),
    ("07", "task_comparison.ipynb",
     "Compare heterogeneity and set sizes, disease vs. chest-pain, side by side",
     "Section 9 — Robustness"),
]

# Which notebook(s) are the hands-on material for a given section divider.
_SECTION_NOTEBOOKS = {
    "5": "00, 01",
    "7": "02, 04",
    "8": "02, 05",
    "9": "03, 04, 06, 07",
}


# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------

def new_slide(bg=WHITE):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    _page_counter["n"] += 1
    return slide


def set_notes(slide, text):
    """Attach a PowerPoint speaker note (Presenter View / View > Notes Page)."""
    if text:
        slide.notes_slide.notes_text_frame.text = text


def _tf(slide, left, top, width, height, anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tf


def _run(p, text, size=18, color=INK, bold=False, italic=False, font=FONT):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def add_footer(slide, kicker=""):
    n = _page_counter["n"]
    tf = _tf(slide, Inches(0.5), SH - Inches(0.45), Inches(6), Inches(0.35))
    p = tf.paragraphs[0]
    _run(p, f"{DECK_TITLE}" + (f"  ·  {kicker}" if kicker else ""), size=9, color=MUTED)
    tf2 = _tf(slide, SW - Inches(1.2), SH - Inches(0.45), Inches(0.8), Inches(0.35))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    _run(p2, str(n), size=9, color=MUTED)


def add_header(slide, kicker, title, title_size=28):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    tf = _tf(slide, Inches(0.7), Inches(0.3), SW - Inches(1.4), Inches(0.35))
    p = tf.paragraphs[0]
    _run(p, kicker.upper(), size=13, color=ACCENT, bold=True)
    tf2 = _tf(slide, Inches(0.7), Inches(0.62), SW - Inches(1.4), Inches(0.9))
    p2 = tf2.paragraphs[0]
    _run(p2, title, size=title_size, color=INK, bold=True)
    return Inches(1.65)  # y-offset where body content should start


def add_image(doc_slide, rel_path, left, top, width=None, height=None):
    full = os.path.join(FIG, rel_path)
    if not os.path.exists(full):
        box = doc_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                                          width or Inches(4), height or Inches(3))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = MUTED
        tf = box.text_frame
        tf.word_wrap = True
        _run(tf.paragraphs[0], f"[missing figure: {rel_path}]", size=12, color=MUTED, italic=True)
        return box
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    return doc_slide.shapes.add_picture(full, left, top, **kwargs)


# ----------------------------------------------------------------------------
# Slide-type builders
# ----------------------------------------------------------------------------

def title_slide(title, subtitle_lines, meta_line, author_lines, venue_line, notes=None):
    slide = new_slide(INK)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.55), SW, Inches(0.06))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    tf = _tf(slide, Inches(0.9), Inches(1.5), SW - Inches(1.8), Inches(1.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, title, size=44, color=WHITE, bold=True)

    tf2 = _tf(slide, Inches(1.2), Inches(2.85), SW - Inches(2.4), Inches(1.4))
    for i, line in enumerate(subtitle_lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _run(p, line, size=18, color=RGBColor(0xd8, 0xe4, 0xf3))

    tf3 = _tf(slide, Inches(1.2), Inches(4.15), SW - Inches(2.4), Inches(0.5))
    p = tf3.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, meta_line, size=13, color=RGBColor(0x9f, 0xb7, 0xd6), italic=True)

    tf4 = _tf(slide, Inches(1.2), Inches(5.3), SW - Inches(2.4), Inches(1.0))
    for i, line in enumerate(author_lines):
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _run(p, line, size=15 if i == 0 else 12.5, color=WHITE, bold=(i == 0))

    tf5 = _tf(slide, Inches(1.2), Inches(6.5), SW - Inches(2.4), Inches(0.5))
    p = tf5.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, venue_line, size=12, color=PLACEHOLDER, italic=True)
    set_notes(slide, notes)
    return slide


def section_divider(number, title, subtitle=None, notes=None):
    slide = new_slide(ACCENT_DARK)
    tf = _tf(slide, Inches(1.0), Inches(2.7), Inches(2.0), Inches(1.6))
    p = tf.paragraphs[0]
    _run(p, number, size=90, color=RGBColor(0x6f, 0x9c, 0xd6), bold=True)

    tf2 = _tf(slide, Inches(3.1), Inches(3.0), SW - Inches(4.2), Inches(1.6))
    p2 = tf2.paragraphs[0]
    _run(p2, title, size=34, color=WHITE, bold=True)
    if subtitle:
        p3 = tf2.add_paragraph()
        _run(p3, subtitle, size=15, color=RGBColor(0xcd, 0xdc, 0xef), italic=True)

    nb = _SECTION_NOTEBOOKS.get(number)
    if nb:
        p4 = tf2.add_paragraph()
        p4.space_before = Pt(14)
        _run(p4, f"HANDS-ON: notebooks {nb}", size=13, color=RGBColor(0xff, 0xd6, 0x7e), bold=True)

    add_footer(slide)
    set_notes(slide, notes)
    return slide


def bullet_slide(kicker, title, items, size=18, spacing=12, note=None, title_size=28, notes=None):
    slide = new_slide()
    top = add_header(slide, kicker, title, title_size=title_size)
    body_h = SH - top - Inches(1.0) - (Inches(0.5) if note else Inches(0))
    tf = _tf(slide, Inches(0.8), top, SW - Inches(1.6), body_h)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        if isinstance(item, tuple):
            text, kind = item
        else:
            text, kind = item, "bullet"
        if kind == "bullet":
            _run(p, "▪  ", size=size, color=ACCENT, bold=True)
            _run(p, text, size=size, color=INK)
        elif kind == "sub":
            p.level = 1
            _run(p, "–  ", size=size - 2, color=MUTED)
            _run(p, text, size=size - 2, color=INK_2)
        elif kind == "head":
            _run(p, text, size=size + 2, color=ACCENT_DARK, bold=True)
    if note:
        tfn = _tf(slide, Inches(0.8), SH - Inches(1.05), SW - Inches(1.6), Inches(0.6))
        _run(tfn.paragraphs[0], note, size=12, color=MUTED, italic=True)
    add_footer(slide, kicker)
    set_notes(slide, notes)
    return slide


def two_col_slide(kicker, title, left_head, left_items, right_head, right_items,
                   size=15.5, title_size=26, notes=None):
    slide = new_slide()
    top = add_header(slide, kicker, title, title_size=title_size)
    col_w = (SW - Inches(2.0)) / 2
    for col, (head, items) in enumerate([(left_head, left_items), (right_head, right_items)]):
        left = Inches(0.7) + col * (col_w + Inches(0.6))
        tf = _tf(slide, left, top, col_w, SH - top - Inches(0.9))
        p = tf.paragraphs[0]
        _run(p, head, size=17, color=ACCENT_DARK, bold=True)
        for item in items:
            pp = tf.add_paragraph()
            pp.space_before = Pt(8)
            _run(pp, "▪  ", size=size, color=ACCENT, bold=True)
            _run(pp, item, size=size, color=INK)
    add_footer(slide, kicker)
    set_notes(slide, notes)
    return slide


def image_slide(kicker, title, rel_path, caption=None, note=None, img_width=None, title_size=26, notes=None):
    slide = new_slide()
    top = add_header(slide, kicker, title, title_size=title_size)
    avail_h = SH - top - Inches(0.55) - (Inches(0.9) if caption else Inches(0))
    w = img_width or Inches(9.6)
    pic = add_image(slide, rel_path, int((SW - w) / 2), top, width=w)
    # re-center vertically within available space and clip height if too tall
    if pic.height > avail_h:
        scale = avail_h / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
        pic.left = int((SW - pic.width) / 2)
    pic.top = int(top)
    cap_top = pic.top + pic.height + Inches(0.12)
    if caption:
        tf = _tf(slide, Inches(0.8), cap_top, SW - Inches(1.6), Inches(0.85))
        _run(tf.paragraphs[0], caption, size=12.5, color=MUTED, italic=True)
    if note:
        tfn = _tf(slide, Inches(0.8), SH - Inches(0.9), SW - Inches(1.6), Inches(0.5))
        _run(tfn.paragraphs[0], note, size=11.5, color=MUTED, italic=True)
    add_footer(slide, kicker)
    set_notes(slide, notes)
    return slide


def two_image_slide(kicker, title, rel_path_a, rel_path_b, caption=None, title_size=26, notes=None):
    slide = new_slide()
    top = add_header(slide, kicker, title, title_size=title_size)
    avail_h = SH - top - Inches(0.55) - (Inches(0.7) if caption else Inches(0))
    half_w = (SW - Inches(2.0)) / 2
    for i, rel in enumerate((rel_path_a, rel_path_b)):
        left = Inches(0.7) + i * (half_w + Inches(0.6))
        pic = add_image(slide, rel, left, top, width=half_w)
        if pic.height > avail_h:
            scale = avail_h / pic.height
            pic.height = int(pic.height * scale)
            pic.width = int(pic.width * scale)
            pic.left = int(left + (half_w - pic.width) / 2)
        pic.top = int(top)
    if caption:
        tf = _tf(slide, Inches(0.8), SH - Inches(0.85), SW - Inches(1.6), Inches(0.65))
        _run(tf.paragraphs[0], caption, size=12.5, color=MUTED, italic=True)
    add_footer(slide, kicker)
    set_notes(slide, notes)
    return slide


def table_slide(kicker, title, headers, rows, col_widths=None, font_size=11,
                 header_size=11.5, note=None, title_size=26, extra_bullets=None, notes=None):
    slide = new_slide()
    top = add_header(slide, kicker, title, title_size=title_size)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    extra_h = Inches(0.0)
    bullets_h = Inches(0)
    if extra_bullets:
        bullets_h = Inches(0.32 * len(extra_bullets) + 0.15)
    table_h = SH - top - Inches(0.55) - bullets_h - (Inches(0.45) if note else Inches(0))
    table_w = SW - Inches(1.4)
    gtable = slide.shapes.add_table(n_rows, n_cols, Inches(0.7), top, table_w, table_h).table

    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            gtable.columns[i].width = int(table_w * (w / total))

    for j, h in enumerate(headers):
        cell = gtable.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        tf.clear()
        _run(tf.paragraphs[0], h, size=header_size, color=WHITE, bold=True)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = gtable.cell(i + 1, j)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if i % 2 == 1 else WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            _run(tf.paragraphs[0], str(val), size=font_size, color=INK)

    y = top + table_h + Inches(0.12)
    if extra_bullets:
        tf = _tf(slide, Inches(0.8), y, SW - Inches(1.6), bullets_h)
        for i, b in enumerate(extra_bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _run(p, "▪  ", size=13.5, color=ACCENT, bold=True)
            _run(p, b, size=13.5, color=INK_2)
        y = y + bullets_h
    if note:
        tf = _tf(slide, Inches(0.8), y, SW - Inches(1.6), Inches(0.4))
        _run(tf.paragraphs[0], note, size=11, color=MUTED, italic=True)
    add_footer(slide, kicker)
    set_notes(slide, notes)
    return slide


def callout_slide(kicker, title, quote, source=None, notes=None):
    slide = new_slide(LIGHT)
    top = add_header(slide, kicker, title)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(2.6),
                                  SW - Inches(2.6), Inches(2.6))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, quote, size=22, color=ACCENT_DARK, bold=True, italic=True)
    if source:
        tf2 = _tf(slide, Inches(1.3), Inches(5.4), SW - Inches(2.6), Inches(0.5))
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        _run(p2, source, size=13, color=MUTED)
    add_footer(slide, kicker)
    set_notes(slide, notes)
    return slide


def closing_slide(lines, repro_cmd=None, notes=None):
    slide = new_slide(INK)
    tf = _tf(slide, Inches(1.2), Inches(2.2), SW - Inches(2.4), Inches(1.2))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, "Thank you", size=44, color=WHITE, bold=True)

    tf2 = _tf(slide, Inches(1.2), Inches(3.4), SW - Inches(2.4), Inches(1.6))
    for i, line in enumerate(lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _run(p, line, size=15, color=RGBColor(0xd8, 0xe4, 0xf3))

    if repro_cmd:
        tf3 = _tf(slide, Inches(1.2), Inches(5.3), SW - Inches(2.4), Inches(1.0))
        p = tf3.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, repro_cmd, size=11, color=RGBColor(0x9f, 0xb7, 0xd6), italic=True, font="Consolas")
    set_notes(slide, notes)
    return slide


# ============================================================================
# BUILD
# ============================================================================

# ---- 1. Title -------------------------------------------------------------
title_slide(
    "Beyond Interoperability",
    ["Hands-On Federated ML for Research Data Curation Infrastructure",
     "Technical Companion: Conformal Prediction for Site-Level Heterogeneity in Federated Clinical Data"],
    "UCI Heart Disease federation  ·  4 hospitals  ·  920 patients",
    ["Arnab Mukherjee", "Oklahoma State University — Center for Health Science, Tulsa"],
    "[Workshop / venue name and date — add here]",
    notes="\"Good morning, everyone, and welcome. My name is Arnab Mukherjee, I'm with the Center for "
          "Health Science at Oklahoma State University, and today we're going to spend the next few "
          "hours going Beyond Interoperability — building, by hand, the tools that tell us whether data "
          "moving between hospitals can actually be trusted once it gets there.\n\n"
          "\"Here's how today is going to work. Roughly half of what we do will be me talking through "
          "ideas like this slide, and the other half will be hands-on — you, in a Jupyter notebook, "
          "running real code on real hospital data. We have eight notebooks, numbered 00 through 07, "
          "and I'll tell you exactly when to open each one as we go.\n\n"
          "\"One thing I want to say clearly up front: you do not need a statistics background or a "
          "machine-learning background to follow today. Every idea — every formula, every acronym — "
          "gets explained before I show you the math. If at any point something "
          "sounds like jargon, stop me and ask.\n\n"
          "\"Before we go any further — everyone, please open your terminal now and make sure you have "
          "the repository cloned and your Python environment set up, exactly as described in the "
          "README. Take two minutes, get that running, and give me a thumbs up when you're ready. "
          "We will not move on until everyone's environment works, because we're going to need it in "
          "about ten minutes.\"",
)

# ---- 2. Abstract ------------------------------------------------------------
bullet_slide("Abstract", "The one idea to remember from today", [
    "Getting data to move between hospitals is a solved problem. Getting that data to mean the same thing once it arrives is not — and that gap is what this whole workshop is about.",
    "The worst kind of failure is a model that looks perfectly healthy on average while quietly getting one hospital's patients wrong — with nothing in its own output warning you it happened.",
    "This workshop uses a tool called conformal prediction to catch that kind of silent failure, and to measure exactly how much two hospitals' data actually differ, before anyone trusts a shared model on both.",
    "The goal: tell apart genuine differences between patient populations from differences that are really just measurement habits — and give a curation team a concrete number they can act on, not a vague worry.",
    "It's hands-on throughout: real Python exercises, code you can reuse, and diagnostics you can point at your own data after today.",
], notes="\n\n"
         "\"There are three ideas we're going to build today. The first is interoperability — the "
         "ability to actually get a patient's data out of one hospital's computer system and into a "
         "form another system can read. The second is federated learning — a way for several hospitals "
         "to train one shared model together without a single patient record ever leaving its own "
         "hospital. And the third is conformal prediction — a way to check, with real numbers, whether "
         "that shared model can actually be trusted once you use it at a hospital it's never seen "
         "before.\n\n"
         "\"Here's why you need all three, not just one. Interoperability gets the data moving, but it "
         "tells you nothing about whether you can trust it once it arrives. Federated learning lets you "
         "train across hospitals without needing to solve that trust question first — but it still "
         "doesn't tell you whether the result is safe to use. Conformal prediction is the piece that "
         "actually answers that — and it only has something to check because the first two ideas got "
         "us to a shared model in the first place. By the end of today, you will have built all three, "
         "yourself, in code.\n\n"
         "\"Say this sentence with me, because we're coming back to it all day: getting data to move "
         "is not the same as being able to trust it once it arrives.\"")

# ---- 3. Agenda --------------------------------------------------------------
two_col_slide("Agenda", "What this deck covers",
    "", ["1. Introduction", "2. Challenges", "3. Motivation",
         "4. The Difficulty of Working With Real Clinical Data",
         "5. Our Dataset: The UCI Heart Disease Federation"],
    "", ["6. Methods of Interoperability", "7. Why Conformal Prediction Helps",
         "8. The Mathematics", "9. Pipeline Walkthrough",
         "References & Appendix"],
    size=18,
    notes="\"Quickly, here's the map. Sections one through four are scene-setting — I'll spend about "
          "twenty minutes convincing you this is a real, live problem, not something invented for a "
          "workshop slide. Section five introduces the dataset we use for everything — four real "
          "hospitals, nine hundred twenty patients. Sections six through eight build up our three core "
          "ideas one at a time: data standards, then conformal prediction, then the math underneath it. "
          "Section nine is where we put every piece together into one working pipeline and try to break "
          "it four different ways.\n\n"
          "\"The next slide is your run sheet for the day — it tells you exactly which notebook to have "
          "open during which section. Keep it visible if you can.\"")

# ---- 3b. Workshop notebook sequence ------------------------------------------
table_slide("Agenda", "Workshop notebook sequence",
    ["#", "Notebook", "What it does", "Runs during"],
    [[n, f, w, s] for n, f, w, s in NOTEBOOKS],
    col_widths=[0.5, 2.6, 5.0, 2.9], font_size=12.5, header_size=13,
    extra_bullets=["Run these in order, 00 → 07 — each takes roughly 10-20 minutes and produces some of the figures shown later in this deck.",
                   "Every notebook is a narrated walk-through of the same reusable code library — nothing is built twice, so what you see live is exactly what generated every slide.",
                   "Two extra clinical questions (Section 9's resting-ECG and exercise-angina results) come from a standalone script rather than a dedicated notebook — flagged again when that part of the deck is reached."],
    notes="\"This table is the map for the whole workshop — eight notebooks, 00 through 07, and I'll tell "
          "you exactly when to open each one. You'll build a conformal predictor on "
          "a simple model first, in about five lines of code, before we wrap that same idea around a "
          "more complicated federated model later, in notebook 04.\n\n"
          "\"One more thing before we start: two of today's later results — the resting-ECG and "
          "exercise-angina findings in Section 9 — come from a standalone script, not a dedicated "
          "notebook. I'm telling you now so nobody goes looking for a notebook 08 that doesn't exist.\"")

# ---- Section 1: Introduction -----------------------------------------------
section_divider("1", "Introduction",
    notes="\"Let's start with the first idea: interoperability. No notebook yet — this section is short, "
          "about five minutes — I just want to ask one question: "
          "data can already move between hospitals today, pretty easily. So why isn't that enough?\"")

bullet_slide("Introduction", "Moving data and trusting data are two different problems", [
    "“Interoperability” means two computer systems can send and read each other's data without a person manually translating in between. That part is largely solved today.",
    "It says nothing about whether the data means the same thing once it lands somewhere new.",
    "A concrete example: two hospitals can both have a “cholesterol” field, in the exact same format, in the exact same spot — and still be describing a different measurement habit, if one hospital tests every patient and the other only tests when there's a specific reason to.",
    "That gap — data successfully arriving vs. data being safe to trust the same way everywhere — is the problem this entire project is built around.",
], notes="\"Interoperability just means two systems can send each other data and both read it, "
         "without a person sitting in the middle translating by hand. Today, for the most part, that's "
         "a solved problem. Hospitals can send each other files, APIs exist, formats are standardized.\n\n"
         "\"But here's what interoperability does not tell you: whether the data means the same thing "
         "once it arrives. Let me give you a concrete example. Two hospitals can both have a field "
         "called 'cholesterol.' Same column name, same units, same format. And they can still mean "
         "something completely different by it — because one hospital tests every single patient's "
         "cholesterol, and the other only orders the test when there's already a reason to suspect a "
         "problem. Same field, different meaning underneath.\n\n"
         "\"Hold onto that example, because it's not hypothetical — in Section 5, I'm going to show you "
         "this exact gap, measured, in our own dataset: one of our four hospitals has zero recorded "
         "cholesterol readings, for every single patient. We'll come back to this.\"")

bullet_slide("Introduction", "Conformal prediction: a tool for catching that gap", [
    "Think of it as a wrapper you can put around any existing model — it turns a single confident guess into a short, honest list of plausible answers.",
    "It comes with a promise: build enough of these lists, and the true answer will be somewhere on the list at least as often as you asked for (say, 90% of the time) — but only measured at the hospital the model was calibrated on.",
    "Take that same model and list-building recipe to a second hospital, and the promise is only as good as the assumption that the two hospitals' patients look statistically alike.",
    "When they don't look alike, the promise breaks — and, crucially, breaks in a way you can measure and point to, instead of failing silently the way a raw confidence score does.",
], notes="\"Here's the clearest way I can put this: conformal prediction is a magnifying glass, not a "
         "better model. It does not make the underlying model any smarter. What it does is take the "
         "model's blind spots — the places where it's wrong but doesn't know it's wrong — and makes "
         "those blind spots visible and measurable instead of invisible. That's the entire pitch of "
         "this workshop's second half, and in Section 7 I'll show you that exact pitch with real numbers "
         "from real hospitals.\"")

bullet_slide("Introduction", "What you'll actually build today", [
    "The real UCI Heart Disease dataset, loaded from its four original hospitals: Cleveland, Hungary, Switzerland, and the V.A. Long Beach — not a cleaned-up, pre-packaged version.",
    "A federated-learning simulator, built from scratch, so you can see exactly how hospitals train one shared model without ever pooling their patients' records.",
    "A conformal-prediction library, also built from scratch, implementing the two list-building recipes used throughout the workshop.",
    "A full set of diagnostics for telling apart “these hospitals' patients are genuinely different” from “these hospitals just measure things differently.”",
    "Dozens of figures, each paired with a plain-language explanation, the math behind it for anyone who wants it, and an honest note on what's solid versus what would need more work.",
], notes="\"Everything you'll touch today is built from scratch, in plain NumPy — no imported "
         "conformal-prediction library, no imported federated-learning library. I did that on purpose, "
         "so that every single line is something you can read and understand, not a black box you have "
         "to trust because a package did it for you.\n\n"
         "\"By the end of this workshop you'll have four real hospitals' data, trained one shared model "
         "across three of them without moving a single patient record, built a conformal predictor "
         "from scratch, and used it to catch that model failing at a fourth hospital it never saw. "
         "That's the whole workshop, in one sentence.\"")

# ---- Section 2: Challenges --------------------------------------------------
section_divider("2", "Challenges",
    notes="\"Still no notebook yet — I want to spend a few minutes making sure you feel, "
          "why this is genuinely important, before we start building things in Section 5. Let's go through "
          "five things that make this problem tricky.\"")

bullet_slide("Challenges", "Five things that make this genuinely hard", [
    "Telling apart a real difference between patient populations from a difference that's just how carefully something was measured — with no outside answer key to check against.",
    "Silent failure: a model's own confidence score gives you no warning at all that it has wandered outside the population it was trained on.",
    "You usually can't just combine everyone's data to sidestep the first two problems — patient records typically have to stay inside their own hospital by law and policy.",
    "Two hospitals using the identical data format doesn't mean they mean the same thing by every field — shared structure isn't shared meaning.",
    "The usual scorecards (overall accuracy, one blended coverage number) hide every one of the above — you need to check results broken out by hospital, by diagnosis, and by patient subgroup to see any of it.",
], size=17.5, notes="\"Let me walk through these five, because every one of them comes back later "
         "today, and I want you to recognize them when they do.\n\n"
         "\"Number one and number two — telling real difference from measurement habit, and a model "
         "that fails silently — those two are exactly what conformal prediction, our third big idea, "
         "is built to solve. Number three — you can't just pool everyone's data — is exactly why "
         "federated learning, our second big idea, exists at all. Number four — same format doesn't "
         "mean same meaning — is what Section 6's tour of data standards is about. And number five — "
         "blended scorecards hide all of this — is why, for the rest of today, I am going to insist "
         "we look at results broken out by hospital, never just one blended average.\n\n"
         "\"Remember these five numbers. We're going to answer every one of them by the end of the workshop.\"")

# ---- Section 3: Motivation --------------------------------------------------
section_divider("3", "Motivation",
    notes="")

two_col_slide("Motivation", "Why this matters right now, not just academically",
    "3.1  AI in medicine is moving faster than cross-hospital checking",
    ["A model that tests well at one hospital regularly does worse at another — different patient mix, different equipment, different measurement habits all play a role.",
     "Regulators are catching up: the FDA's approach to AI-based medical software now expects tracking of real-world performance after deployment, broken out by patient subgroup — not just one number measured before launch."],
    "3.2  Federated learning has gone from research idea to practical necessity",
    ["The core method used in this workshop (FedAvg) was introduced by McMahan et al. in 2017, and quickly became well known through its use in mobile phone keyboard prediction.",
     "In healthcare specifically, gathering every hospital's patient records into one place is often simply not legally possible — so this approach is closer to the only option than a stylistic preference.",
     "A well-known 2022 benchmark collection (FLamby) uses this exact same four-hospital dataset as its own standard test case — this is a recognized, real testbed, not an invented toy problem."],
    size=15,
    notes="\"Two things are happening in parallel right now, out in the real world. First, AI models in "
          "medicine are being deployed faster than anyone is checking whether they hold up across "
          "hospitals — a model that looks great at the hospital where it was built regularly does worse "
          "somewhere else, and regulators have noticed: the FDA now expects tracking of real-world "
          "performance after a model launches, not just one clean number measured before.\n\n"
          "\"Second — federated learning, our second big idea, has gone from a research curiosity to "
          "something close to a necessity. The method we're using today, called FedAvg, was introduced "
          "back in 2017, and you're going to build it yourself, from scratch, in notebook 03. In "
          "healthcare specifically, pooling every hospital's patient records into one "
          "place is often just not legally possible choice. And I'll add one more thing: a well-known 2022 benchmark collection called "
          "FLamby uses this exact same four-hospital dataset as its own standard test case.")

two_col_slide("Motivation", "Rules require the data to move; they don't require it to mean the same thing",
    "3.3  Current data-sharing rules mandate that data CAN move, not that it means the same thing everywhere",
    ["In the US, federal rules now require standardized APIs for health data exchange, and a nationwide exchange network is already operating.",
     "The EU has its own parallel regulation (the European Health Data Space), phasing in from 2027.",
     "Both kinds of rules guarantee the bytes can travel between systems — neither one guarantees that a value means the same clinical fact once it arrives somewhere new."],
    "3.4  A checkable uncertainty number fills exactly that gap",
    ["Conformal prediction's promise doesn't depend on assuming any particular shape for the data — it works essentially the same way regardless of what the underlying numbers look like.",
     "It doesn't require the model to be any particular type, or built any particular way.",
     "It hands a reviewer a number they can go check for themselves empirically — not something they have to simply take on faith."],
    size=15,
    notes="I want you to walk away with from this slide: today's regulations "
          "guarantee that health data CAN move between systems. Not one of them guarantees that the "
          "data means the same clinical fact once it gets there. That gap — the difference between "
          "'can move' and 'can be trusted' — is exactly what our third idea, conformal prediction, "
          "exists to fill. It doesn't care what the data looks like, it doesn't care what kind of model "
          "you used, and it hands a reviewer a number they can go check for themselves, not something "
          "they have to take on faith. Sections 7 through 9 are entirely about that number.\"")

callout_slide("Motivation", "The gap this workshop is built to close",
    "Hospitals today can share data with each other. What's still missing is a checkable way to know whether that shared data can actually be trusted the same way once it arrives — this workshop's toolkit exists to make that trust something you can measure, not just assume.",
    notes="…\n\n"
          "\"Read that slide with me: hospitals today can share data with each other. What's still "
          "missing is a checkable way to know whether that shared data can actually be trusted the "
          "same way once it arrives. That's the whole workshop, in one sentence.\"\n\n")

# ---- Section 4: The Difficulty of Real Clinical Data ------------------------
section_divider("4", "The Difficulty of Working With Real Clinical Data",
    notes="")

table_slide("Real Clinical Data", "Stacking up the barriers to a public download",
    ["Barrier", "What it requires", "What that means in practice"],
    [
        ["Ethics review", "Formal review or a documented exemption for nearly any use of patient records",
         "Weeks to months, repeated at every hospital involved"],
        ["Privacy protections", "Careful, expert-verified removal of identifying details",
         "Requires dedicated specialized expertise"],
        ["Data-sharing agreements", "A negotiated legal contract between every pair of institutions sharing data",
         "Gets harder roughly by the square of how many hospitals are involved"],
        ["Access credentialing", "Signed agreement plus required training, per individual researcher",
         "Days to weeks, and has to be renewed periodically"],
        ["Net effect", "Most clinical data never becomes a public, freely reusable research dataset",
         "Federated learning becomes close to a practical necessity, not a preference"],
    ], col_widths=[2.4, 4.4, 3.2], font_size=13, header_size=13,
    notes="\"This table is just the last slide's argument, stacked up as hurdles. Read the last row "
          "with me: because of everything above it, most clinical data never becomes a public, "
          "reusable research dataset. That's the whole reason federated learning matters — it's not an "
          "academic exercise, it's often the only door that's actually open.\"")

# ---- Section 5: Our Dataset --------------------------------------------------
section_divider("5", "Our Dataset", "The UCI Heart Disease Federation",
    notes="\"Alright — this is where we get hands-on. Everyone, please open "
          "notebooks/00_input_data_exploration.ipynb now. Take a second... good.\n\n"
          "\"Run the first cell — it just imports our toolkit and loads all four hospitals' data into "
          "memory. Give it a moment to finish, you should see it print the patient counts... Everyone "
          "have output? Good, let's keep going.\"\n\n"
          "[Notebook 01, site_heterogeneity, comes in partway through this section — I'll tell you "
          "exactly when to switch, once we've talked about the two kinds of hospital-to-hospital "
          "difference.]")

bullet_slide("Our Dataset", "Why this particular dataset", [
    "Four real, independent hospitals all filled out the exact same intake form — so the data genuinely can move between them cleanly.",
    "And yet, underneath that shared form, the four hospitals' patients and practices differ in ways that genuinely matter.",
    "Every difference shown in this section is a real effect actually present in the data — nothing here was constructed or exaggerated to make a point.",
    "It carries none of the access barriers just described in Section 4 — open, free to use, and properly citable.",
], notes="\"Four reasons I picked this dataset for today. One: four real, independent hospitals filled "
         "out the exact same intake form, so the data genuinely does move cleanly between them — that's "
         "our first idea, interoperability, already satisfied. Two: underneath that shared form, these "
         "four hospitals' patients and practices differ in ways that actually matter — that's the "
         "heterogeneity we're about to go find. Three: every difference I show you in this section is a "
         "real effect, sitting in the data right now — I haven't constructed or exaggerated anything to "
         "make a point. And four: zero access barriers, thanks to that historical accident I just "
         "mentioned. Let's go find these differences ourselves.\"")

table_slide("Our Dataset", "Variable dictionary (UCI Heart Disease, ID 45)",
    ["Variable", "Role", "Type", "Description", "Units", "Missing"],
    [
        ["age", "Feature", "Integer", "Age of the patient", "years", "no"],
        ["sex", "Feature", "Categorical", "Sex of the patient", "–", "no"],
        ["cp", "Feature", "Categorical", "Chest pain type", "–", "no"],
        ["trestbps", "Feature", "Integer", "Resting blood pressure (on admission)", "mm Hg", "no"],
        ["chol", "Feature", "Integer", "Serum cholestoral", "mg/dl", "no"],
        ["fbs", "Feature", "Categorical", "Fasting blood sugar > 120 mg/dl", "–", "no"],
        ["restecg", "Feature", "Categorical", "Resting electrocardiographic results", "–", "no"],
        ["thalach", "Feature", "Integer", "Maximum heart rate achieved", "–", "no"],
        ["exang", "Feature", "Categorical", "Exercise induced angina", "–", "no"],
        ["oldpeak", "Feature", "Integer", "ST depression induced by exercise", "–", "no"],
        ["slope", "Feature", "Categorical", "Slope of the peak exercise ST segment", "–", "no"],
        ["ca", "Feature", "Integer", "Major vessels (0-3) by flourosopy", "–", "yes"],
        ["thal", "Feature", "Categorical", "Thalassemia stress-test result", "–", "yes"],
        ["num", "Target", "Integer", "Diagnosis of heart disease (severity)", "–", "no"],
    ], col_widths=[1.1, 0.9, 1.1, 3.6, 0.9, 1.0], font_size=11, header_size=12,
    note="Source: UCI Machine Learning Repository, archive.ics.uci.edu/dataset/45/heart+disease",
    notes="\"This is a reference table — I'm not going to read every row, just skim it with me. Notice "
          "only two fields, ca and thal, are flagged missing here. Keep that in mind, because a couple "
          "of slides from now I'm going to show you a much bigger missingness problem — cholesterol, "
          "completely unrecorded at one entire hospital — and it won't show up in this table at all, "
          "because it's stored as a placeholder value, not a blank. Lesson for today: always look at "
          "what your data actually contains, not just what the documentation claims it contains.\"")

table_slide("Our Dataset", "Coded sub-values behind every categorical field",
    ["Variable", "Coded sub-values"],
    [
        ["sex", "1 = male;  0 = female"],
        ["cp", "1 typical angina;  2 atypical angina;  3 non-anginal pain;  4 asymptomatic"],
        ["fbs", "1 = fasting blood sugar > 120 mg/dl;  0 = false"],
        ["restecg", "0 normal;  1 ST-T wave abnormality;  2 probable/definite LV hypertrophy"],
        ["exang", "1 = exercise induced angina;  0 = no"],
        ["slope", "1 upsloping;  2 flat;  3 downsloping"],
        ["ca", "0, 1, 2 or 3 major vessels colored by flourosopy"],
        ["thal", "3 normal;  6 fixed defect;  7 reversable defect"],
        ["num (target)", "0 no disease;  1–4 increasing severity / vessels affected"],
    ], col_widths=[2.2, 8.4], font_size=13.5, header_size=14,
    notes="\"Also reference material — skim it, don't read it line by line. One thing worth flagging: "
          "one single patient in this dataset has a vessel count value outside this documented 0-to-3 "
          "range entirely. Almost certainly a decades-old data-entry slip, and it's been left in the "
          "data exactly as found. File that away — it's another small example of 'read what's actually "
          "there, not just what the codebook promises.'\"")

table_slide("Our Dataset", "The four hospitals", ["Hospital", "Institution", "Patients", "Share with any heart disease", "Most common severity"],
    [
        ["Cleveland", "Cleveland Clinic Foundation", "303", "45.9%", "No disease (54.1%)"],
        ["Hungary", "Hungarian Inst. of Cardiology", "294", "36.1%", "No disease (63.9%)"],
        ["Switzerland", "Univ. Hospital Zurich / Basel", "123", "93.5%", "Mild (39.0%)"],
        ["V.A. Long Beach", "V.A. Medical Center", "200", "74.5%", "Mild (28.0%)"],
        ["All four combined", "–", "920", "55.3% overall", "–"],
    ], col_widths=[1.6, 3.2, 1.2, 2.0, 2.0], font_size=14, header_size=14,
    extra_bullets=["The share of patients with any heart disease nearly triples across these four hospitals, from 36% to 94% — that's a genuine difference in who each hospital sees, not a data-entry problem."],
    notes="\"Everyone, run the site-summary cell in notebook 00 right now, the one that follows this "
          "point in the notebook — you should see this exact table print out on your own screen, "
          "computed live from the four original hospital files. Don't just take my slide's word for "
          "it — watch it come out of your own code.\"")

image_slide("Our Dataset", "Patient counts and disease rates, by hospital", "01_site_overview.png",
    caption="Disease prevalence swings from 36% at Hungary to 94% at Switzerland — a real difference in which patients each hospital sees, not a data-entry mistake.",
    img_width=Inches(8.6),
    notes="\"This is notebook 00's first real figure — everyone should see this render on your own "
          "screen too. Look at that swing: 36% at Hungary, 94% at Switzerland. Hold that number in your "
          "head, because it's one half of a two-part story I'm about to tell you. This one is a "
          "genuinely different group of patients. The next slide shows you the other half — a "
          "difference that's not about the patients at all, just about what each hospital wrote down.\"")

table_slide("Our Dataset", "Two very different kinds of hospital-to-hospital difference", ["Measurement", "Cleveland", "Hungary", "Switzerland", "V.A.", "What kind of difference is this?"],
    [
        ["Cholesterol never recorded", "0.0%", "0.0%", "100.0%", "24.5%", "Measurement habit (specific to one hospital)"],
        ["Vessel count never recorded", "1.3%", "98.6%", "95.9%", "99.0%", "Measurement habit (almost everywhere except Cleveland)"],
        ["Thalassemia test never recorded", "0.7%", "90.5%", "42.3%", "83.0%", "Measurement habit (almost everywhere except Cleveland)"],
        ["How severe the disease is", "mixed", "skews healthy", "skews sicker", "mixed", "Real difference in who the hospital sees"],
    ], col_widths=[1.9, 1.2, 1.2, 1.4, 1.0, 3.3], font_size=12, header_size=13,
    notes="\"This table is the heart of the whole section, so let's slow down. Bottom row: a real "
          "difference between patient populations — that's biology. Top three rows: just differences "
          "in what each hospital happened to write down — that's measurement habit, not biology. "
          "Almost everything else in this section is evidence for one of these rows or the other.\n\n")

two_image_slide("Our Dataset", "How much is missing, and what cholesterol looks like when it is recorded",
    "02_missingness.png", "03_chol_distributions.png",
    caption="Left: the fraction of patients missing each measurement, by hospital — cholesterol's complete 100% gap at Switzerland stands out, alongside the much bigger, near-universal gap in vessel count and thalassemia results outside Cleveland. Right: cholesterol values where they were actually recorded (Switzerland has none to show). Why is Switzerland at 100%? Every one of its patients has cholesterol coded as exactly 0 mg/dl in the original source file — a value that's biologically impossible for a living patient, so it's clearly a placeholder for “not measured,” not a real reading. This dataset was assembled from four separate hospital studies run independently in 1988–89, each using its own case-report form; the most likely explanation is that serum cholesterol simply wasn't part of the Zurich/Basel site's own data-collection protocol the way it was at Cleveland — a difference in what got tested, not a data-loss accident. No surviving documentation from the original study states this explicitly, so treat it as the best-supported reading of the evidence, not a confirmed fact.",
    notes="\"Look at Switzerland's cholesterol bar — completely full, 100%. Every single patient. Let's "
          "stop and ask why, because it's worth understanding, not just noting.\n\n"
          "\"Here's what we know for certain: in the original source file, every Swiss patient's "
          "cholesterol is recorded as exactly zero. A living patient cannot have a cholesterol reading "
          "of zero — that's biologically impossible. So this isn't a real measurement, it's a stand-in "
          "for 'we don't have this value.'\n\n"
          "\"Now, why would an entire hospital have zero cholesterol readings? Here's the most likely "
          "explanation, and I want to be honest that this is our best reading of the evidence, not "
          "something a document confirms outright: this dataset wasn't collected as one unified study. "
          "It's four separate hospital studies, run independently back in 1988 and '89, each using its "
          "own case-report form. Cholesterol was very likely just never part of the Zurich and Basel "
          "site's own testing protocol the way it was at Cleveland's. Not lost data. Not a mistake in "
          "typing. A difference, from day one, in what each hospital's own protocol chose to measure.\n\n"
          "\"That's exactly the kind of measurement-habit gap our whole workshop is built to catch — "
          "and notebook 01, which we're about to open, is where you'll compute this exact percentage "
          "yourselves.\"")

table_slide("Our Dataset", "A real difference in outcomes: how severe cases are, by hospital", ["Severity", "Cleveland", "Hungary", "Switzerland", "V.A.", "All combined"],
    [
        ["No disease", "164", "188", "8", "51", "411 (44.7%)"],
        ["Mild", "55", "37", "48", "56", "196 (21.3%)"],
        ["Moderate", "36", "26", "32", "41", "135 (14.7%)"],
        ["Severe", "35", "28", "30", "42", "135 (14.7%)"],
        ["Critical", "13", "15", "5", "10", "43 (4.7%)"],
    ], col_widths=[1.8, 1.3, 1.3, 1.5, 1.1, 1.6], font_size=13.5, header_size=14,
    extra_bullets=["Hungary is mostly “no disease”; Switzerland is almost entirely disease-positive."],
    notes="\"Look at Hungary and Switzerland specifically: 64% of Hungary's patients have no disease at "
          "all; at Switzerland it's just 6.5%. That's genuine biology — real differences in which "
          "patients get referred to each hospital — not a coding mistake. Compare that to the "
          "measurement-habit rows we saw two slides ago — this is the other kind of difference.\"")

two_image_slide("Our Dataset", "The severity split, overall and by hospital",
    "eda/e01_target_distribution.png", "eda/e02_target_by_site.png",
    caption="Left: the overall severity split across all 920 patients, plus a simplified any-disease-or-not summary. Right: the same breakdown, now split out by hospital — the previous table, drawn as a picture.",
    notes="\"Same story as the last table, now as a picture. Quick question before we move on — which "
          "hospital would you guess is going to be hardest for one shared model to generalize to? Take "
          "a guess. ... Hold that guess — Section 7 is going to give you a very concrete answer.\"")

table_slide("Our Dataset", "A simple test: can you guess the hospital just from the patient's numbers?", ["", "Cleveland", "Hungary", "Switzerland", "V.A."],
    [
        ["Cleveland", "50 / 50", "89%", "100%", "88%"],
        ["Hungary", "89%", "50 / 50", "100%", "93%"],
        ["Switzerland", "100%", "100%", "50 / 50", "88%"],
        ["V.A.", "88%", "93%", "88%", "50 / 50"],
    ], col_widths=[1.6, 1.6, 1.6, 1.6, 1.6], font_size=15, header_size=15,
    extra_bullets=["Train a simple classifier whose only job is to guess which hospital a patient came from, using nothing but their clinical measurements — no diagnosis, no label. Every off-diagonal number here is how often it guesses right for that pair. 88%+ everywhere means: yes, easily."],
    notes="\"Here's a test I like: train the simplest possible model, and give it one job — guess which "
          "hospital a patient came from, using nothing but their vital signs and lab values. No "
          "diagnosis, no label, nothing about disease at all. Look at these numbers — 88% or better, "
          "for every single pair of hospitals. That's about as concrete a definition of 'these are "
          "different patient populations' as you'll ever get, and it doesn't even need to know who's "
          "sick.\"")

two_image_slide("Our Dataset", "The same finding, confirmed a second, independent way",
    "05_domain_auc.png", "04_js_divergence.png",
    caption="Left: the previous table's guess-the-hospital accuracy, drawn as a heatmap. Right: a completely different, classifier-free way of measuring how differently shaped two hospitals' data is — applied here to maximum heart rate — pointing to the same conclusion.",
    notes="\"I want to show you the same conclusion reached two completely different ways. On the left, "
          "a trained model guessing the hospital. On the right, no model at all — just directly "
          "comparing the shape of the raw numbers, for maximum heart rate, hospital by hospital. Two "
          "independent methods, same answer. That's what rules out 'maybe this is just one particular "
          "classifier being quirky.' This is notebook 01 — you're computing both of these yourselves "
          "right now.\"")

image_slide("Our Dataset", "Which measurements matter most for telling severity apart", "eda/e03_continuous_grid.png",
    caption="Each panel is one continuous measurement, with five curves (one per severity level, lightest = no disease, darkest = critical). Where the curves pull apart from each other, that measurement is genuinely useful for predicting severity; where they sit on top of each other, it isn't.",
    notes="\"Look at these panels with me and tell me — which ones look most spread apart? ... Right, "
          "maximum heart rate and ST depression tend to separate the most cleanly. Those are our most "
          "informative continuous measurements, and we'll see them again.\"")

image_slide("Our Dataset", "One measurement, isolated by hospital instead of by severity", "eda/e04_thalach_by_site.png",
    caption="Maximum heart rate stood out on the previous slide as one of the more informative measurements for severity. Here it's the same measurement, but grouped by hospital instead — a cleaner look at how its typical range shifts from one hospital to the next.",
    img_width=Inches(8.4),
    notes="\"Same measurement as the last slide, sliced the other way — by hospital instead of by "
          "severity. I want you to notice something important: both kinds of variation are happening "
          "at once in this data. It varies by outcome, and separately, it varies by hospital. Neither "
          "one explains the other away — they're both real, and they're both here.\"")

image_slide("Our Dataset", "The rest of the measurements, by severity", "eda/e05_categorical_grid.png",
    caption="Average severity within each category, for all eight non-numeric measurements — chest-pain type, sex, blood sugar, exercise-induced chest pain, resting-ECG result, ST-segment slope, vessel count, and thalassemia test result. Several show a clear pattern: asymptomatic chest pain, exercise-induced chest pain, an abnormal ST slope, high blood sugar, more affected vessels, and an abnormal thalassemia result all line up with higher average severity.",
    img_width=Inches(9.2),
    notes="\"One panel worth a closer look: vessel count. Average severity climbs smoothly and steadily "
          "as the vessel count goes from 0 up to 3 — the cleanest 'more of this, worse the disease' "
          "pattern on this whole slide. That matters twice over, because vessel count is also the "
          "measurement with the worst missingness outside Cleveland — so it's reassuring to see that, "
          "where we do have it, it behaves exactly the way real cardiology knowledge says it should.\"")

image_slide("Our Dataset", "How every measurement relates to severity, at a glance", "eda/e06_correlation.png",
    caption="How strongly every pair of measurements (and the severity target) move together. Chest-pain type, exercise-induced chest pain, ST depression, maximum heart rate, and vessel count are the strongest, most reliable signals.",
    img_width=Inches(7.2),
    notes="\"This confirms nothing we've seen in the last three slides was a fluke. Five measurements "
          "keep coming up as the strongest signals — chest-pain type, exercise-induced chest pain, ST "
          "depression, maximum heart rate, vessel count. Remember those five names — they'll reappear "
          "for the rest of today.\"")

image_slide("Our Dataset", "The same story, on a completely different question", "cp/c01_class_distribution.png",
    caption="Instead of disease severity, this predicts chest-pain type — a different clinical question, same four hospitals. Hungary sees mostly atypical chest pain; Switzerland is roughly 80% asymptomatic. The same hospital-to-hospital differences show up again here, confirming this isn't specific to one particular question.",
    img_width=Inches(8.6),
    notes="\"Let's pressure-test everything we've just found. Instead of disease severity, this slide "
          "predicts something completely different — chest-pain type — same four hospitals. And look: "
          "the same kind of hospital-to-hospital swing shows up again. That tells us this isn't a quirk "
          "of choosing 'disease severity' specifically — it's a property of these hospitals themselves. "
          "We'll use this exact chest-pain question again later, in Sections 8 and 9.\"")

bullet_slide("Our Dataset", "Recap: what makes this dataset unusually good for teaching this", [
    "Genuinely shareable — all four hospitals used the exact same 13-measurement intake form.",
    "Genuinely different underneath, in two distinguishable ways: real differences between patient populations (disease rates ranging 36%-94%) and simple measurement-habit gaps (cholesterol, and much more severely, vessel count and thalassemia results).",
    "Small enough to train and calibrate live during a single workshop (920 patients total), yet large enough that every effect shown is statistically real, not noise.",
    "Open and free to use — no legal agreement stands between anyone in this room and running the code today.",
    "Every finding shown here reappears on a second, independently chosen clinical question from the same four hospitals.",
], notes="\"Let's recap Section 5 before we move on. This dataset is genuinely shareable, genuinely "
         "different underneath in two distinguishable ways, small enough to run live today but large "
         "enough to be real, completely open, and everything we found replicates on a second clinical "
         "question. Good checkpoint here.\n\n"
         "[If running behind schedule: skip ahead to Section 7 now and only come back to Section 6 if "
         "time allows — Section 6 is background, not something attendees build by hand.]\"")

# ---- Section 6: Methods of Interoperability ----------------------------------
section_divider("6", "Methods of Interoperability", "Choosing a Standard",
    notes="\"Quick detour before we get to the heart of the workshop. No notebook for this section — "
          "about ten minutes of background on data standards, then we're back to hands-on. If we're "
          "running behind, this is the section I'll compress.\"")

bullet_slide("Interoperability", "“Interoperable” isn't just one thing — it comes in degrees", [
    "Think of it as four separate levels: can the data move at all? Is its format standardized? Does a given value mean the same clinical fact everywhere? And is there shared agreement on the policies and workflows around it?",
    "Most of today's regulations do a good job on the first two levels — getting data to move, in a standard format.",
    "The third level — a value actually meaning the same thing everywhere — is where this entire project's heterogeneity findings live. It's the hardest level to mandate by law, and the easiest one to think you've solved when you've really only gotten partway there.",
], notes="\"Our first idea, interoperability, isn't just one switch that's either on or off — it comes "
         "in four levels. Can the data move at all? Is the format standardized? Does a value mean the "
         "same clinical fact everywhere? And is there shared agreement on the policy and workflow "
         "around it? [Sketch these four on a board if you have one.]\n\n"
         "\"Most regulations today are strong on the first two levels. The third level — meaning the "
         "same thing everywhere — is where every finding from Section 5 lives. It's the hardest one to "
         "legislate, and the easiest one to think you've already solved when you've really only gotten "
         "partway there.\"")

table_slide("Interoperability", "A quick tour of today's data-sharing standards",
    ["Standard", "What it's mainly used for", "How well it fits multi-hospital research"],
    [
        ["HL7 v2 messaging", "Older-generation, point-to-point clinical messaging", "Weak fit — built to run hospital operations, not pooled research"],
        ["HL7 FHIR", "Modern web-API-style data exchange (what current mandates require)", "Good for moving data live; not itself a research-analysis format"],
        ["OMOP (OHDSI)", "Standardizing observational data specifically across institutions", "Excellent — built specifically for this exact use case"],
        ["openEHR", "Detailed, template-based clinical record modeling", "Good in principle; not widely adopted by US hospitals yet"],
        ["CDISC SDTM/ADaM", "Formal clinical-trial data submission to regulators", "Weak fit — built for trials, not everyday hospital records"],
        ["DICOM", "Medical imaging exchange and storage", "Not applicable to this kind of tabular data — but the right choice once imaging enters the picture"],
    ], col_widths=[2.2, 4.4, 4.8], font_size=13, header_size=13,
    notes="\"Quick tour, not a deep dive. The one to remember is OMOP — it's genuinely built for "
          "exactly this use case, standardizing observational data across institutions. [If anyone in "
          "the room has hands-on OMOP experience, ask them how it handles the measurement-habit "
          "differences we saw in Section 5.]\"")

bullet_slide("Interoperability", "What a shared format can fix, and what it can't", [
    "It reduces measurement-habit differences between hospitals, but it doesn't eliminate them.",
    "A hospital can code cholesterol using the exact correct standard term and still simply order that test far less often than another hospital does — no format fixes that.",
    "No shared format changes the fact that a specialist referral center sees a sicker, pre-selected group of patients to begin with — that's about who walks in the door, not how the data is written down.",
    "Checking for these gaps has to happen after adopting a shared format, as an ongoing check — not as a one-time substitute for doing that checking at all.",
    "Worth remembering: this project's own dataset predates every standard in the table above — an honest, unfiltered look at what hospital data looks like before any of this standardization happened.",
], size=17, notes="\"Here's the line I want you to remember from this whole section: a hospital can "
         "code cholesterol using the exact correct standard term, and still simply order that test less "
         "often than another hospital. No data standard fixes that. That's precisely why the checks "
         "we're about to build in Sections 7 and 8 have to sit downstream of adopting a shared format, "
         "as an ongoing habit — not as a one-time replacement for it.\n\n"
         "\"Alright — background's done. Let's get to the real center of today: why we need conformal "
         "prediction at all.\"")

# ---- Section 7: Why Conformal Prediction Helps -------------------------------
section_divider("7", "Why Conformal Prediction Helps", "With Site-Level Heterogeneity",
    notes="\"Everyone, pause what you're doing for a second — this is the center of the whole workshop, "
          "so I want your full attention here. Everything we've done up to now has been about noticing "
          "that hospitals differ. This section is about the tool that catches a model quietly failing "
          "because of it — that's conformal prediction, our third big idea.\n\n"
          "\"In a moment I'll ask you to open notebook 02, conformal_basics — that's where you'll build "
          "a conformal predictor yourselves, in about five lines of code. Later in this section we'll "
          "come back to notebook 04, conformal_under_site_shift, for the headline result of the entire "
          "day. Stay with me for the next few slides first.\"")

callout_slide("Conformal Prediction", "An everyday version of the problem",
    "A weather forecaster who says “90% chance of rain” should be right about 9 days out of 10 whenever they say that. If they're only right 6 times out of 10, you stop trusting the number — even though every individual forecast still sounded confident.",
    source="A model that reports “92% confident” should work the same way. Conformal prediction is the tool that checks whether it actually does — and says so honestly when it doesn't.",
    notes="\"Before I show you a single number from our dataset, let me ask you something. Has anyone "
          "here ever stopped trusting a weather app, because it kept saying '90% chance of rain' and "
          "that just didn't feel right anymore? ... [Wait for a show of hands or a nod.] That instinct "
          "— checking a confident-sounding number against what actually happened — is exactly what "
          "conformal prediction does automatically, for a clinical model instead of a weather app. A "
          "forecaster who says 90% chance of rain should be right about nine days out of ten, whenever "
          "they say that number. If they're only right six times out of ten, you stop trusting the "
          "number, even though every individual forecast still sounded just as confident. A model that "
          "says 92% confident should work exactly the same way — and conformal prediction is the tool "
          "that checks whether it actually does, honestly, out loud, instead of us just hoping.\"")

bullet_slide("Conformal Prediction", "The problem: a model can sound confident and still be wrong", [
    "A model says “92% chance of heart disease.” That number is not automatically trustworthy — nothing forces the model to actually be right 92% of the time.",
    "Move that same model to a new hospital, and it often keeps sounding just as confident even while its real, checkable accuracy quietly drops.",
    "The model has no way to know it has left home. It can't tell you it's now guessing in unfamiliar territory — that silence is the dangerous part.",
    "Getting the data to move between hospitals (interoperability) is a separate problem from whether the model's judgment still holds once it arrives. Solving the first does nothing for the second.",
], notes="\"Here's the single most important idea I want you to leave with today: confident, but wrong. "
         "A model says 92% confident. That number sounds precise, sounds trustworthy. But nothing "
         "forces it to actually be right 92% of the time — it's just what the model happens to output.\n\n"
         "\"Let me ask the room: has anyone here taken a model, or honestly even just a rule of thumb, "
         "that worked well in one place, and used it somewhere new without double-checking it still "
         "applied? ... That's this exact failure mode. And here's the dangerous part — the model itself "
         "gives you no warning sign when it happens. It doesn't get quieter. It doesn't hedge. It just "
         "keeps sounding exactly as confident as before, while quietly being wrong more often.\n\n"
         "\"And notice — this is a completely different problem from interoperability. Getting the data "
         "to move between hospitals is solved. Whether the model's judgment still holds once that data "
         "arrives is not. Solving the first does nothing for the second.\"")

bullet_slide("Conformal Prediction", "What conformal prediction actually gives you", [
    "Instead of one confident guess, it hands back a short list of plausible answers — for example, “the true severity is somewhere in {1, 2, 3}.”",
    "That list comes with a promise attached: statements built this way are right at least 90% of the time, if you ask for 90%. Not a hope — a guarantee you can go check.",
    "The guarantee holds no matter how good or bad the underlying model is — the only thing it needs is that the patients used to set the threshold look statistically like the patients it's later asked to judge.",
    "When a new hospital's patients don't look like that anymore, the promise breaks — and because we can measure the promise, we can catch the break instead of it happening silently.",
], notes="\"So here's the fix. Instead of one confident guess, conformal prediction hands back a short "
         "list of plausible answers — 'the true severity is somewhere in one, two, or three' — and that "
         "list comes with an actual promise attached: built this way, it's right at least 90% of the "
         "time, if you asked for 90%. Not a hope. A guarantee you can go check, empirically, with real "
         "data.\n\n"
         "\"Everyone, open notebook 02, conformal_basics, now. Run the cells with me — you're about to "
         "build exactly this kind of list, from scratch, in about five lines of NumPy. Go ahead... "
         "Watch what comes out — that's a real, checkable prediction set, built by code you just ran "
         "yourself. We'll wrap this same idea around our federated model later, in notebook 04.\n\n"
         "\"One more thing before we move on: the guarantee doesn't need the model to be any good. It "
         "only needs one thing — that the patients you set the threshold on look statistically like the "
         "patients you later use it on. When a new hospital's patients stop looking like that, the "
         "promise breaks. And because we can measure the promise, we catch that break instead of it "
         "happening silently.\"")

image_slide("Conformal Prediction", "Two ways an average can lie to you", "paper/fig10_coverage_notions.png",
    caption="Picture two patient groups. Left: no promise at all — the model could be right any fraction of the time. Middle: the promise holds only on average — the two groups' results are blended together, so one can quietly fail while the blended number still looks fine. Right: the promise holds for every group separately — nothing is hidden by blending. Recreation of Angelopoulos & Bates (2022), Figure 10.",
    img_width=Inches(8.0),
    notes="\"Look at these three pictures — this is a cartoon version of a real result I'm about to "
         "show you, so let's understand it here first, before the real numbers land. Left: no promise "
         "at all. Middle: the promise holds only on average — the two groups get blended together, so "
         "one group can quietly fail while the blended number still looks perfectly healthy. Right: the "
         "promise holds for every group separately, nothing hidden by blending.\n\n"
         "\"Say this with me: 'meets the target on average' and 'meets the target everywhere' are two "
         "different claims — and only the second one is actually safe to act on. Keep that in your head "
         "for the next slide.\"")

image_slide("Conformal Prediction", "The headline result: watch the promise break", "09_coverage_by_site.png",
    caption="A conformal predictor is calibrated using only Cleveland's patients, then deployed at all four hospitals unchanged. It keeps its 90% promise at Cleveland, Hungary and the V.A. — and silently drops to 71% at Switzerland. The shaded band marks ordinary random wobble; 71% falls well outside it, so this is a real effect, not noise.",
    img_width=Inches(8.4),
    notes="\"Everyone, this is it — this is the single most important figure in the entire workshop. "
          "Everything we've done all day has been building to this moment, so let's take our time.\n\n"
          "\"If you have notebook 04 open, run the headline cell now and watch the numbers print with "
          "me live. Here's what we did: we set our threshold using only Cleveland's patients — nothing "
          "else. Then we deployed that exact same model and threshold at all four hospitals, including "
          "Switzerland, which this model has never seen in training or calibration.\n\n"
          "\"Watch the numbers: Cleveland, 90.8%. Hungary, 90.8%. The V.A., 90.5%. And Switzerland... "
          "70.7%. Look at that gap. The shaded band on this chart is ordinary random wobble — and 71% "
          "sits well outside it. This is not noise. This is a real, measurable failure, and we caught "
          "it because we checked.\"")

table_slide("Conformal Prediction", "Why you have to check every site, not just the average", [
        "Deployment site", "Empirical coverage", "Gap vs. 90% promise", "Verdict"],
    [
        ["Cleveland (calibration site)", "90.8%", "+0.8 pp", "Promise kept"],
        ["Hungary", "90.8%", "+0.8 pp", "Promise kept"],
        ["V.A. Long Beach", "90.5%", "+0.5 pp", "Promise kept"],
        ["Switzerland (held out entirely)", "70.7%", "−19.3 pp", "Promise broken, silently"],
    ], col_widths=[3.6, 2.4, 2.4, 2.4], font_size=15, header_size=15,
    extra_bullets=["Calibrated only on Cleveland's 303 patients. Blend all four sites into one overall coverage number and it lands close to 90% — looking perfectly healthy. Only by checking Switzerland on its own does the 19-point failure become visible."],
    notes="\"I want to say this slowly, because it's the practical lesson you should take back to your "
          "own work: if you blended all four hospitals into one overall coverage number, it would land "
          "close to 90% and look perfectly healthy. It would completely hide Switzerland's failure. The "
          "only way to catch it is to check every hospital on its own, never just the average. That "
          "check is cheap — you can bolt it onto any model already running in production, starting "
          "today.\"")

image_slide("Conformal Prediction", "Every hospital pair, tested at once", "08_transfer_matrix.png",
    caption="Each row is 'the hospital we set the threshold using'; each column is 'the hospital we then deploy at.' Reading straight down the diagonal (same hospital both times) is always healthy. Anything set using Cleveland, Hungary or the V.A. and then deployed at Switzerland lands around 70-72%; a threshold set using Switzerland instead over-delivers everywhere else.",
    img_width=Inches(7.6),
    notes="\"Let's not just take my word for the Cleveland case — let's check all sixteen possible "
          "pairings at once. Each row is 'which hospital we set the threshold with,' each column is "
          "'which hospital we deployed at.' Straight down the diagonal is always healthy, that's using "
          "a hospital on itself.\n\n"
          "\"Now look at the Switzerland column — anything set using Cleveland, Hungary, or the V.A. "
          "lands around 70 to 72% there. And here's a twist worth noticing: flip it around, set the "
          "threshold using Switzerland, and it doesn't fail elsewhere — it over-delivers, because "
          "Switzerland is the sickest population, so its threshold is looser everywhere else. These "
          "mismatches don't always fail in the scary direction, but they're still a mismatch worth "
          "catching every time.\"")

# ---- Section 8: The Mathematics ----------------------------------------------
section_divider("8", "The Mathematics", "And What Every Piece of Code Produces",
    notes="\"Now let's look under the hood. I want to reassure you before we start: you do not need a "
          "statistics background for this section. Every formula I show you gets translated into a "
          "plain sentence right next to it — I just want you to recognize the same three steps every "
          "single time: score something, find a cutoff, draw a boundary. That's it, that's the whole "
          "section.\n\n"
          "\"This is the densest part of the day, so we'll take it slowly, and there are a few slides "
          "near the end marked optional — totally fine to skip those if we're short on time. Notebook "
          "02, revisited, and notebook 05, recreating_paper_figures, are our hands-on material here.\"")

image_slide("Mathematics", "What a prediction set actually looks like", "paper/fig01_prediction_set_examples.png",
    caption="Three real patients from this project's data, all truly \"No disease,\" and the list of plausible answers the calibrated model actually returns for each — a confident single answer, a so-so 3-answer list, and a genuinely unsure 4-answer list. In that hardest case, the model's own best guess is wrong — yet the correct answer, \"No disease,\" still made the list. That is the whole point of the guarantee: even when the model is confused, the true answer still shows up somewhere in what it hands you.",
    img_width=Inches(9.4),
    notes="\"Before I show you a single formula, look at these three real patients — not made-up "
          "examples, three actual people from our dataset, all of whom truly have no disease. Look at "
          "the third one, the hardest case: the model's own top guess is actually wrong there. And yet "
          "'no disease,' the true answer, still made the list it returned.\n\n"
          "\"That's not a coincidence — that's the entire point of the guarantee. Every formula I show "
          "you for the rest of this section is just 'how do we build a list with that exact property.' "
          "Keep this patient in mind as we go.\"")

bullet_slide("Mathematics", "The one assumption everything rests on", [
    "In plain terms: the patients used to set the threshold have to look statistically like the patients the model will later be asked to judge. Statisticians call this “exchangeability” — think of it as “no surprises between calibration and deployment.”",
    "It does not require the data to follow a bell curve, or be cleanly separable, or come from any particular kind of model — about as weak and general an assumption as a guarantee can rest on.",
    "This is exactly the assumption cross-hospital deployment can break outright: Cleveland's patients (46% have any disease) and Switzerland's patients (94% have any disease) simply don't look alike statistically.",
    "This dataset was chosen specifically because it breaks that assumption in a real, measurable way — so the failure is something you can see and quantify, not just describe.",
], notes="\"Everything we've built rests on exactly one assumption. Statisticians call it "
         "'exchangeability,' but don't worry about that word — here's the plain version: the patients "
         "you set your threshold on have to look statistically like the patients you'll later use it "
         "on. No surprises between calibration and deployment. That's it. It doesn't need bell-curve "
         "data, doesn't need any particular kind of model — about as weak an assumption as you can ask "
         "for, which is exactly why it works with anything.\n\n"
         "\"And here's the punch line: this is exactly the assumption that moving a model between two "
         "very different hospitals can break outright. Cleveland's patients: 46% have any disease. "
         "Switzerland's: 94%. Those two groups simply don't look statistically alike — and I picked "
         "this dataset specifically because it breaks that assumption in a way we can measure, not "
         "just describe.\"")

bullet_slide("Mathematics", "The recipe, in three steps: score, find a cutoff, draw the boundary", [
    "Step 1 — score: for every calibration patient, ask “how surprising was their true answer to the model?” (a nonconformity score, written s(x, y)). A confident, correct answer scores low; a confident, wrong answer scores high.",
    "Step 2 — find a cutoff: sort those scores and take a percentile near the top, called q̂. Formally,  q̂ = Quantile(calibration scores;  ⌈(n + 1)(1 − α)⌉ / n) — the ⌈⌉ and the “+1” are just a small correction so the guarantee is exact even with a limited number of patients, not just true on average over many repeats.",
    "Step 3 — draw the boundary: for a new patient, keep every possible answer whose surprise score doesn't clear that cutoff:  C(x) = { y : s(x, y) ≤ q̂ }.",
    "The guarantee this produces:  P(true answer is in the returned list) ≥ 1 − α — true for any scoring rule, and any underlying model, confident or not.",
    "Two scoring rules used in this toolkit: LAC keeps things simple (score = 1 − the model's probability on the true class) and tends to return the shortest lists. APS is a bit more adaptive — it adds up sorted probabilities until they cover the true class — and tends to size each list to how genuinely uncertain that particular patient is.",
], size=15.5, notes="\"Three steps. Say them with me: score, cutoff, boundary. Step one, score — for "
         "every calibration patient, ask how surprising their true answer was to the model. Confident "
         "and correct scores low, confident and wrong scores high. Step two, find a cutoff — sort all "
         "those scores and take a percentile near the top, we call it q-hat. That little ceiling symbol "
         "and the plus-one in the formula are just a small correction so the promise is exact even with "
         "a few hundred patients, not just true on average over infinite repeats. Step three, draw the "
         "boundary — for a new patient, keep every answer whose surprise score doesn't clear that "
         "cutoff.\n\n"
         "\"That's the whole recipe. Every remaining formula today is a variation on these same three "
         "steps. This exact cutoff calculation is what you ran yourselves in notebook 02 — if you want, "
         "open that cell again and step through it line by line with me now.\"")

two_image_slide("Mathematics", "The two scoring rules, worked through visually",
    "paper/fig02_conformal_illustration.png", "paper/fig04_adaptive_prediction_sets.png",
    caption="Left (LAC): score a calibration patient → plot all the calibration scores as a histogram and mark the cutoff q̂ → for a new patient, keep every answer whose probability clears that bar. Right (APS): same three steps, but answers are sorted from most to least likely and added up until the running total crosses q̂.",
    notes="\"Here are our two scoring rules, side by side. LAC on the left, APS on the right — same "
          "three-step recipe we just walked through, just a different way of measuring 'how surprising "
          "was this.' LAC is simpler and tends to give shorter lists on average. APS is a bit more "
          "adaptive — we use it for the multiclass secondary tasks in Section 9 because it sizes each "
          "patient's list more sensibly to how uncertain that specific patient actually is.\"")

table_slide("Mathematics", "LAC vs. APS, head to head on the same model", [
        "Scoring rule", "Average list length", "Coverage achieved"],
    [
        ["LAC (simpler rule)", "3.09 out of 5 answers", "90.3% (target 90%)"],
        ["APS (adaptive rule)", "2.96 out of 5 answers", "90.3% (target 90%)"],
    ], col_widths=[3.6, 3.6, 3.2], font_size=15.5, header_size=15.5,
    extra_bullets=["Same model, same calibration data, same 90% target — only the scoring rule changes. Both keep their promise; APS does it with slightly shorter lists on average, because it adapts list length to how uncertain each individual patient is rather than using one fixed rule for everyone."],
    notes="\"If you still have notebook 02 open, run both cells back to back — the LAC one, then the "
          "APS one — and watch what prints. Same model, same calibration patients, same 90% target, "
          "only the scoring rule changes. Both hit 90.3% coverage. But look at the list lengths — APS "
          "does it with slightly shorter lists on average. Here's the takeaway: the scoring rule you "
          "pick doesn't change whether the promise holds, only how efficiently it's kept.\"")

image_slide("Mathematics", "Not a toy example — this project's own numbers", "07_calibration_scores.png",
    caption="Every one of these bars is a real calibration score from this project's own federated model, trained across Cleveland, Hungary and the V.A. The red line is the cutoff q̂ — everything to its left is what makes the 90% promise hold.",
    img_width=Inches(8.4),
    notes="\"I want to draw a sharp contrast with the last slide's clean, textbook illustrations — this "
          "histogram is messier, because it's real. Every bar here is an actual calibration score from "
          "our own federated model. The red line is the cutoff q-hat that goes on to produce Section "
          "7's headline result. Not a toy example — this is genuinely the number our pipeline "
          "computed.\"")

image_slide("Mathematics", "How many patients do you need before you can trust the number?", "11_coverage_beta.png",
    caption="Run the calibration step many times on small random samples and the achieved coverage bounces around — sometimes a little high, sometimes a little low, just from the luck of which patients happened to land in the sample. At 50 patients that bounce is wide; at 1,000 it's tight. Switzerland's 71% sits far outside even the widest realistic bounce at Cleveland's own scale (303 patients) — so it isn't bad luck, it's a genuine site effect.",
    img_width=Inches(8.2),
    notes="\"A sharp question some of you are probably already asking: is Switzerland's 71% real, or "
          "just noise from a smallish calibration set? Let's answer that directly. Run the calibration "
          "step over and over on small random samples, and the achieved coverage bounces around a bit "
          "— just from the luck of the draw. At 50 patients, that bounce is wide. At 1,000, it's tight. "
          "Now — even generously, at Cleveland's own scale of 303 patients, 71% sits far outside any "
          "realistic bounce. That's how we know Switzerland's result is a genuine site effect, not "
          "sampling noise. [This comes from a nearly century-old statistical result — Vovk, 2012 — but "
          "nobody needs the name to follow the argument, just the picture.]\"")

image_slide("Mathematics", "The same reassurance, at a much bigger scale", "paper/fig11_coverage_distribution.png",
    caption="The same bounce-from-sample-size idea, now shown at the much larger scales used in the original paper (100, 1,000, and 10,000 calibration patients) — the spread narrows sharply as the sample grows. Recreation of Angelopoulos & Bates (2022), Figure 11.",
    img_width=Inches(8.4),
    notes="\"Same idea as the last slide, just zoomed out to a much bigger scale — 100, 1,000, and "
          "10,000 calibration patients. Watch how sharply that spread narrows. [If anyone's silently "
          "thinking '303 patients is already plenty,' this is the slide that answers them — at these "
          "bigger scales, the bounce basically disappears.]\"")

bullet_slide("Mathematics", "Coverage alone isn't the whole story", [
    "There's a cheap, useless way to hit 90% coverage: always return every possible answer. Technically correct, completely unhelpful to a clinician.",
    "A genuinely useful list is short when the model is confident and long when it's genuinely unsure — and it should visibly grow at hospitals where that particular patient's outcome is harder to call.",
    "One overall coverage number can also hide a second, subtler failure: it can be right on average across all patients while being badly wrong for one specific true diagnosis.",
    "Concrete example: on the chest-pain task, the rarest true diagnosis (typical angina, only 23 patients) is only covered 26% of the time against a 90% promise — a failure that's completely invisible in the overall pooled number.",
], notes="\"Quick thought experiment: what's the cheapest possible way to hit 90% coverage, guaranteed? "
         "... Just always return every possible answer. Technically correct. Completely useless to a "
         "clinician. So coverage by itself isn't the whole story — a genuinely useful list is short "
         "when the model's confident and long when it's genuinely unsure.\n\n"
         "\"There's a second, subtler failure too: one overall coverage number can be right on average "
         "while being badly wrong for one specific diagnosis. Concrete example, from our own chest-pain "
         "results — the rarest true diagnosis, typical angina, only 23 patients — is covered just 26% "
         "of the time against a 90% promise. Invisible in the overall number. Let's look at both of "
         "these side by side on the next slide.\"")

two_image_slide("Mathematics", "Two more things worth checking, side by side",
    "10_set_sizes.png", "cp/c04_class_conditional.png",
    caption="Left: how long the returned lists are, by hospital (disease task) — longer lists at a hospital mean the model is systematically less sure about that hospital's patients. Right: coverage broken out by true chest-pain diagnosis — the rarest diagnosis (typical angina, 23 patients) is covered only 26% of the time, a failure the overall 90% number never reveals.",
    notes="\"Two independent failure modes, deliberately side by side. On the left: how long the lists "
          "are, by hospital — longer lists at a hospital mean the model itself is less sure about that "
          "hospital's patients, visible even before you check coverage at all. On the right: that 26% "
          "number from the last slide, now as a picture — coverage broken out by true diagnosis. Both "
          "of these matter directly if your pipeline serves smaller, under-represented patient groups.\"")

bullet_slide("Mathematics", "The same trick works well beyond classification", [
    "The three-step recipe from earlier — score how surprising something is, find a cutoff, use the cutoff to draw a boundary — isn't specific to picking from a list of diagnoses.",
    "Predicting a number instead of a category (e.g. a lab value)? Take a fitted range and widen it by exactly the cutoff amount to reach the same guarantee.",
    "Have a single best-guess number plus some notion of how uncertain it is? Widen a band around that guess the same way.",
    "Working with a full Bayesian model instead? The same idea turns its output into a region above a probability threshold.",
    "These next three slides use small made-up examples, exactly as the original paper does — they're here to show the idea generalizes, not as results from this project's own clinical data.",
], notes="\"Quick note before the next three slides: they're optional depth, using small made-up "
         "examples exactly like the original paper does, not run on our own patients. Perfectly fine to "
         "skip these if we're tight on time — I'm telling you now so nobody thinks there's a missing "
         "notebook they were supposed to run. The point of these three is just: this recipe — score, "
         "cutoff, boundary — isn't specific to picking a diagnosis from a list. It works for a numeric "
         "prediction, for a single best-guess number, and even inside a full Bayesian model.\"")

two_image_slide("Mathematics", "The same trick, applied to numbers instead of categories",
    "paper/fig06_conformalized_quantile_regression.png", "paper/fig08_uncertainty_scalar.png",
    caption="Left: a fitted range of likely numeric values (dashed) widened by the cutoff amount to reach the guarantee (solid). Right: a single best-guess number with a symmetric widened band around it. Both are made-up illustrative examples, recreating Angelopoulos & Bates (2022), Figures 6 and 8.",
    notes="\"Same 'widen by the cutoff amount' idea, now applied outside classification entirely. Left, "
          "a numeric range. Right, a single best-guess number with a band around it. [Most relevant if "
          "anyone in the room mainly works with numbers — lab values, dosages, wait times — rather than "
          "categories.]\"")

image_slide("Mathematics", "The same trick, inside a Bayesian model", "paper/fig09_conformalized_bayes.png",
    caption="Instead of a category list or a numeric range, the returned answer is a region: every outcome whose estimated probability clears a threshold. Made-up illustrative example, recreating Angelopoulos & Bates (2022), Figure 9.",
    img_width=Inches(7.0),
    notes="\"Last of our three generalization examples, and a good closing point for this whole section: "
          "score, cutoff, boundary really is one general-purpose trick. It's not something invented "
          "just for picking between diagnoses — it works here too, inside a full Bayesian model.\"")

# ---- Section 9: Pipeline Walkthrough ------------------------------------------
section_divider("9", "Pipeline Walkthrough", "and Justification: Is This Enough?",
    notes="\"Alright — we've built all three ideas separately: interoperability got us a shared, "
          "comparable dataset; conformal prediction gave us a way to check trust; and in a moment "
          "we're going to build federated learning too, the piece that lets hospitals train together "
          "without pooling data. This section puts all three together into one real pipeline, start to "
          "finish, and then we stress-test it four different ways to make sure our headline finding "
          "wasn't a fluke of one particular dataset choice.\n\n"
          "\"Notebook 03, federated_training, then notebook 04 again, revisited, are our core hands-on "
          "material here. Notebooks 06 and 07 cover the robustness checks later in this section.\"")

bullet_slide("Pipeline — Training", "How the hospitals train one shared model together", [
    "One shared model learns to weigh each of the 13 patient measurements and combine them into a severity guess — nothing fancier than a weighted combination, about 70 numbers total, small enough to print out and read.",
    "Each round, the server sends the current shared model out to the three training hospitals — Cleveland, Hungary, the V.A. Patient records themselves never leave any hospital; only the model's numbers travel.",
    "Each hospital nudges its own copy slightly, using only its own patients, to make it fit its local data a bit better.",
    "The server then averages the three nudged copies back into one, weighting each hospital by how many patients it contributed — a bigger hospital's nudge counts for more.",
    "That whole cycle repeats 200 times. Switzerland is deliberately left out of training entirely, so it can later serve as a genuinely unseen hospital.",
], size=16.5, notes="\"Everyone, open notebook 03, federated_training, now. This is our second big "
         "idea, built from scratch, live. Let's read the first cell together before running it.\n\n"
         "\"Here's what's about to happen: one shared model — about 70 numbers total, small enough to "
         "print out and actually read — gets sent out to three training hospitals: Cleveland, Hungary, "
         "the V.A. Patient records themselves never leave their hospital. Only those 70 numbers travel. "
         "Each hospital nudges its own copy slightly, using only its own patients. Then the server "
         "averages the three nudged copies back into one, weighted by how many patients each hospital "
         "contributed.\n\n"
         "\"Go ahead and run the training cell now — it repeats that whole cycle 200 times. While it "
         "runs, notice: Switzerland is deliberately left out of training entirely. Hold onto that, "
         "because it's about to matter a lot.\"")

image_slide("Pipeline — Training", "Confirming the training actually finished learning", "06_fed_curves.png",
    caption="Each line is one hospital's error rate as training progresses. All three are still dropping visibly early on and go essentially flat by round ~150 — a sign that training ran long enough to settle, rather than being cut off partway through.",
    img_width=Inches(8.4),
    notes="\"Your training cell should have produced a chart like this one — let's check it together. "
          "Each line is one hospital's error rate. Notice they're all still dropping early on, and go "
          "essentially flat around round 150. That flat tail is what tells us training ran long enough "
          "to actually settle, rather than being cut off partway through. If your curves look like "
          "this, training worked.\"")

image_slide("Pipeline — Training", "Same check, a completely different prediction target", "cp/c02_fed_curves.png",
    caption="The identical training setup, now predicting chest-pain type instead of disease severity. Same settling-down pattern — a sign this behavior is a property of the method, not a coincidence specific to one target.",
    img_width=Inches(8.4),
    notes="\"Same exact check, but for a completely different question — chest-pain type instead of "
          "disease severity. This comes from notebook 06, later in the day. Same settling-down pattern. "
          "That's how we know this is a property of the method itself, not a coincidence specific to "
          "one target.\"")

table_slide("Pipeline — Training", "Does sharing models beat pooling all the data?", [
        "Hospital", "Shared-model accuracy", "Hypothetically-pooled accuracy"],
    [
        ["Cleveland", "60.7%", "61.1%"],
        ["Hungary", "71.1%", "70.7%"],
        ["Switzerland (never trained on)", "21.1%", "20.3%"],
        ["V.A. Long Beach", "44.0%", "44.0%"],
    ], col_widths=[3.4, 3.4, 3.6], font_size=15, header_size=15,
    extra_bullets=["“Pooled” means a hypothetical model trained as if every patient record could sit in one place — the comparison this project is not allowed to run in a real hospital setting, only here as a reference point. The shared model trained without ever pooling data lands within a point of it everywhere — the privacy-preserving approach costs almost nothing in accuracy."],
    notes="\"Let's answer the question everyone's probably wondering: does this privacy-preserving "
          "approach actually cost us anything? Look at this table with me — 'pooled' here means the "
          "hypothetical version where every patient record sits in one place, which we're never "
          "actually allowed to do in a real hospital setting. Compare the two columns: Cleveland, "
          "60.7% versus 61.1%. Hungary, 71.1% versus 70.7%. The V.A., 44.0% versus 44.0%, identical. "
          "Federated training, without ever pooling a single patient record, lands within about a "
          "point of the pooled version everywhere. Say this plainly to the room: federated learning "
          "isn't a compromise here — it's close to a free win, and patient data never has to move.\n\n"
          "\"One clarification before we continue: Switzerland's low accuracy here, around 21%, is a "
          "completely separate story from the coverage result we're about to see. Accuracy and "
          "coverage are different questions — accuracy is 'did it get the diagnosis right,' coverage "
          "is 'did the true answer make the list.' We're about to focus entirely on the second one.\"")

bullet_slide("Pipeline — Training", "Under the hood: how each hospital's nudge is computed", [
    "In plain terms: after each hospital scores its own patients, it looks at how far off those scores were and adjusts its copy of the model a small step in the direction that would have made them less wrong. Repeated a few times, that's \"training.\"",
    "The formal version, for anyone who wants it: the loss is cross-entropy plus a small penalty that discourages overly large weights (L(θ) = −mean(log P[i, yᵢ]) + (λ/2)·‖W‖², λ = 1e-2, applied only to the weights, not the bias).",
    "Weight gradient:  ∇W = Xᵀ(P − Y) / n + λW  — “P minus Y” is just “what the model guessed minus what actually happened,” averaged over patients.",
    "Bias gradient:  ∇b = mean(P − Y, axis=0)  — the same idea, just for the model's baseline offset.",
    "Update rule — the plainest possible version, no shortcuts or acceleration tricks:  θ ← θ − lr·∇L(θ), a fixed step size applied a few times per round at every hospital, so every line of it can be read and hand-verified rather than trusted as a black box.",
], size=15.5, notes="\"In plain terms, here's what a 'nudge' actually is: each hospital looks at how "
         "far off its guesses were, and adjusts its copy of the model a small step in the direction "
         "that would have made it less wrong. Do that a few times, that's training. For anyone who "
         "wants the exact formulas, they're on the slide, and every single line is something you can "
         "open and read in notebook 03. If we have time, let's open src/fedconformal/federated.py "
         "together and match these formulas to the actual code.\"")

bullet_slide("Pipeline — Calibration", "Setting the threshold that makes the promise exact", [
    "Once the shared training rounds finish, the model's numbers are frozen — nothing about the model itself changes again from this point on.",
    "Each training hospital's patients are then split in half: one half sets the threshold (calibration), the other half is kept aside purely to double-check it later — neither half was used to train the model in the first place.",
    "The frozen model scores the calibration half only; those scores are what Section 8's “how surprising was this” recipe runs on.",
    "The cutoff picked from those scores becomes q̂ — the one and only number this whole process still has to “fit” once the model itself is done training.",
    "That cutoff is only trustworthy for patients who look statistically like the ones it was set on — which is exactly the assumption the next slide puts to the test.",
], size=16, notes="\"Training is done, model's frozen — nothing about it changes from here on. Now we "
         "add our third idea, conformal prediction, on top. Each training hospital's patients get "
         "split in half — one half sets the threshold, the other half is set aside purely to "
         "double-check it later. Neither half touched training. That cutoff — q-hat — is the one "
         "number left to 'fit' once the model itself is done. And it's only trustworthy for patients "
         "who look statistically like the ones it was set on. Let's go put that exact claim to the "
         "test.\"")

bullet_slide("Pipeline — Deployment", "The real test: an entirely unfamiliar hospital", [
    "The toughest version of the test: set the threshold using only Cleveland's patients (303 of them) — nothing from any other hospital.",
    "Then deploy that frozen model and fixed threshold at every hospital, including Switzerland — a hospital the model never trained on and the threshold never saw.",
    "No retraining, no adjusting, no sneaking a look at Switzerland's answers first — just running the frozen model forward and applying the same fixed threshold everywhere.",
    "The result: the 90% promise holds at Cleveland, Hungary and the V.A. — and quietly slips to 70.7% at Switzerland (see the table earlier in this section). The promise breaks exactly where Switzerland's patients stop looking statistically like Cleveland's.",
], size=16.5, notes="\"This is it — everyone, open notebook 04 again if you closed it, and find the "
         "final cell. This is the moment the entire workshop has been building to, so let's slow all "
         "the way down.\n\n"
         "\"Here's the setup: set the threshold using only Cleveland's 303 patients — nothing else. "
         "Deploy that frozen model, unchanged, at every hospital, including Switzerland, which it has "
         "never trained on and never calibrated against. No retraining. No adjusting. No peeking at "
         "Switzerland's true answers first.\n\n"
         "\"Run that final cell now, together with me, and watch the numbers print one hospital at a "
         "time. Cleveland... Hungary... the V.A.... and Switzerland. There it is — 70.7%. You just "
         "watched, with your own code, a model's promise silently break exactly where the patients "
         "stopped looking statistically alike. That's the whole workshop, in one cell.\"")

table_slide("Pipeline — Robustness", "Is this specific to one clinical question?", [
        "What we're measuring", "Disease severity (5 possible answers)", "Chest-pain type (4 possible answers)"],
    [
        ["How differently the answer is distributed across hospitals", "0.128 (bigger gap)", "0.073"],
        ["How easily a hospital can be told apart from its patients alone", "0.930", "0.929 (essentially the same)"],
        ["Average promise kept across all four hospitals (target 90%)", "90.9%", "91.0%"],
        ["Worst single hospital's promise kept", "87%", "88%"],
        ["Average list length returned", "2.68 out of 5", "2.25 out of 4"],
    ], col_widths=[5.0, 2.8, 2.8], font_size=13.5, header_size=13.5,
    extra_bullets=["Same full pipeline, run start to finish a second time on a completely different clinical question. Disease severity shows a bigger gap across hospitals — no surprise, since its prevalence swings further (36%–94%) than chest-pain type's does — but both questions keep every hospital close to the 90% promise."],
    notes="\"Fair question to ask at this point: did we just get lucky, picking disease severity as our "
          "one example? Let's check, using notebook 07. This runs the entire pipeline a second time, "
          "start to finish, on a completely different question — chest-pain type. Look at these two "
          "columns: nearly identical how-separable-are-the-hospitals score, about 0.93 both times, "
          "despite very different amounts of label imbalance. Say this plainly: that similarity tells "
          "us we're looking at a property of these hospitals themselves, not an accident of which "
          "clinical question we happened to ask.\"")

image_slide("Pipeline — Robustness", "The same comparison, as a picture", "compare/task_comparison.png",
    caption="The five numbers from the previous table, side by side, disease-severity vs. chest-pain-type. The same story shows up in both: real cross-hospital differences, but the 90% promise holds close to target everywhere.",
    img_width=Inches(9.6),
    notes="\"Same five numbers, now as a picture. Before I move on, guess with "
          "me: which of the two questions do you think shows the bigger cross-hospital gap? ... Right, "
          "disease severity, because its prevalence swings further across our four hospitals.\"")

two_image_slide("Pipeline — Chest-Pain Task", "Running the full headline test a second time", "cp/c05_transfer_matrix.png", "cp/c06_coverage_by_site.png",
    caption="Left: every hospital-pair combination for the chest-pain question — the same all-pairs test as Section 7's disease-severity version, just applied to a different diagnosis. Right: a Cleveland-only threshold, tested at all four hospitals.",
    notes="\"If you want to see the full parallel rather than just trust my summary table, this is "
          "notebook 06's own repeat of Section 7's two headline figures — same exact tests, different "
          "clinical question.\"")

two_image_slide("Pipeline — Chest-Pain Task", "Calibration scores and per-diagnosis coverage",
    "cp/c03_calibration_scores.png", "cp/c04_class_conditional.png",
    caption="Left: the chest-pain question's own calibration-score histogram — the same kind of chart shown earlier for disease severity. Right: coverage broken out by true chest-pain diagnosis — the rarest one (typical angina) is covered just 26% of the time.",
    notes="\"Here's that 26% number from Section 8, back in its original context — coverage on the rare "
          "'typical angina' diagnosis, only 23 patients. Remember this from earlier? Same finding, "
          "shown here where it actually comes from.\"")

image_slide("Pipeline — Chest-Pain Task", "How long are the returned lists, by hospital?", "cp/c07_set_sizes.png",
    caption="The chest-pain question's version of the list-length-by-hospital chart shown earlier for disease severity.",
    img_width=Inches(8.0),
    notes="\"Rounds out the chest-pain question's full diagnostic set — same list-length chart we saw "
          "for disease severity back in Section 8.\"")

bullet_slide("Pipeline — Further Targets", "Pushing further: two more clinical questions", [
    "Two questions (disease severity, chest-pain type) already show this isn't a fluke of one particular label — but could it still be a fluke of those two specific labels?",
    "A companion script reruns the exact same federated-plus-conformal pipeline on two more columns from the same four hospitals: the resting-ECG result (three possible readings) and whether the patient has exercise-induced chest pain (yes/no).",
    "Both were picked because they're almost fully filled in across every hospital (missing for only 2 patients and 55 patients respectively out of 920) — unlike a few other columns that are missing for the large majority of patients at some hospitals and would leave a question barely answerable there.",
    "The yes/no question runs through a simpler two-outcome version of the model automatically — but every downstream check (scoring rule, coverage, list length) works exactly the same way regardless.",
], size=16.5, notes="\"Let's push one step further. Two questions already show this isn't a fluke of "
         "one particular label — but could it still be a fluke of those two specific labels? A "
         "companion script reruns our exact pipeline on two more columns: resting-ECG result, and "
         "whether the patient has exercise-induced chest pain. I picked these two specifically because "
         "they're almost completely filled in at every hospital — unlike some other columns that are "
         "missing for most patients at some hospitals, which would leave a question barely answerable "
         "there.\n\n"
         "\"One flag: this pair runs from a standalone script, not a notebook — the one part of today "
         "without its own notebook. I'm telling you so nobody goes looking for it.\"")

table_slide("Pipeline — Further Targets", "The worst-hospital promise, across four independent questions", [
        "Clinical question", "Possible answers", "Worst single hospital", "Gap vs. 90% promise"],
    [
        ["Disease severity", "5", "70.7%", "−19.3 pp"],
        ["Resting-ECG result", "3", "74.5%", "−15.5 pp"],
        ["Exercise-induced chest pain", "2", "82.8%", "−7.2 pp"],
        ["Chest-pain type", "4", "89.0%", "−1.0 pp"],
    ], col_widths=[3.4, 2.0, 2.6, 2.4], font_size=15, header_size=15,
    extra_bullets=["Three of these four completely independent clinical questions show a real, double-digit-point promise failure at their hardest hospital. The size of each failure tracks how differently that specific measurement is distributed across hospitals — this isn't one generic side effect of combining these four hospitals, it's a real, question-by-question signal."],
    notes="\"Quick question for the room before I reveal this table: of these four completely "
          "independent clinical questions, which one would you guess fails worst at its hardest "
          "hospital? ... Take a guess. [Reveal:] Disease severity, at 70.7%. But look — three of the "
          "four show a real, double-digit-point failure. The size of each failure tracks how "
          "differently that specific measurement is distributed by hospital — this isn't one generic "
          "side effect of combining these four hospitals, it's a real, question-by-question signal.\"")

two_image_slide("Pipeline — Resting-ECG Task", "Third question: resting-ECG result", "restecg/01_class_distribution.png", "restecg/02_fed_curves.png",
    caption="Left: how the resting-ECG reading is distributed by hospital — normal readings dominate at Hungary and Switzerland, while the V.A. leans toward one specific abnormality. Right: the training-settled-down check, same as before, for this question.",
    notes="\"This third question, resting ECG, turns out to be the closest analogue to our original "
          "disease-severity failure — 74.5% worst-hospital coverage, not far behind the 70.7% headline "
          "number. Keep that in mind, it lands in two slides.\"")

two_image_slide("Pipeline — Resting-ECG Task", "Resting-ECG: calibration scores and per-diagnosis coverage",
    "restecg/03_calibration_scores.png", "restecg/04_class_conditional.png",
    caption="Left: the resting-ECG question's own calibration-score histogram. Right: coverage broken out by true resting-ECG reading.",
    notes="\"Same diagnostic battery we've run twice already, now applied to resting ECG — for anyone "
          "who wants to check every angle rather than just trust the headline table.\"")

two_image_slide("Pipeline — Resting-ECG Task", "Resting-ECG: every hospital pair, and the headline test",
    "restecg/05_transfer_matrix.png", "restecg/06_coverage_by_site.png",
    caption="Left: every hospital-pair combination for resting ECG — the V.A. is the hardest hospital to deploy at no matter which hospital set the threshold, landing 74-88% against the 90% promise. Right: a Cleveland-only threshold, tested at all four hospitals.",
    notes="\"Notice something here — the hardest hospital this time is the V.A., not Switzerland. "
          "Worth saying explicitly: the 'hardest hospital' isn't always the same one. It depends "
          "entirely on which specific measurement you're predicting.\"")

image_slide("Pipeline — Resting-ECG Task", "How long are the returned lists, by hospital?", "restecg/07_set_sizes.png",
    img_width=Inches(8.0),
    notes="\"Completes the resting-ECG diagnostic set — same list-length chart as our earlier two "
          "questions.\"")

two_image_slide("Pipeline — Exercise-Angina Task", "Fourth question: exercise-induced chest pain", "exang/01_class_distribution.png", "exang/02_fed_curves.png",
    caption="Left: how common exercise-induced chest pain is by hospital — roughly doubling from Hungary (30%) to the V.A. (65%), a milder version of the same kind of hospital-to-hospital swing seen for disease severity. Right: the training-settled-down check for this question.",
    notes="\"Our fourth and final question — exercise-induced chest pain, a simple yes-or-no answer. It "
          "automatically routes through a simpler two-outcome version of the model, but every check "
          "downstream — scoring rule, coverage, list length — works exactly the same as the other "
          "three questions.\"")

two_image_slide("Pipeline — Exercise-Angina Task", "Exercise angina: calibration scores and per-diagnosis coverage",
    "exang/03_calibration_scores.png", "exang/04_class_conditional.png",
    caption="Left: the exercise-angina question's own calibration-score histogram. Right: coverage broken out by true yes/no answer.",
    notes="\"Same diagnostic battery, one more time, applied to this yes/no question.\"")

two_image_slide("Pipeline — Exercise-Angina Task", "Exercise angina: every hospital pair, and the headline test",
    "exang/05_transfer_matrix.png", "exang/06_coverage_by_site.png",
    caption="Left: every hospital-pair combination for exercise angina — the weakest single pairing (threshold set at Hungary, deployed at Switzerland) lands at 83%, a real gap but a smaller one than resting ECG or disease severity showed. Right: a Cleveland-only threshold, tested at all four hospitals.",
    notes="\"The mildest of our four questions — 82.8% worst-hospital coverage. Still a genuine 7-point "
          "gap below our 90% promise, but noticeably smaller than resting ECG or disease severity.\"")

image_slide("Pipeline — Exercise-Angina Task", "How long are the returned lists, by hospital?", "exang/07_set_sizes.png",
    img_width=Inches(8.0),
    notes="\"That completes our fourth and final robustness check. Let's zoom all the way back out for "
          "a second: four completely different clinical questions, same pipeline, same basic story "
          "every single time. That's not a coincidence — that's the finding.\"")

bullet_slide("Pipeline", "Where this stands, and what's honestly still missing", [
    "For what it set out to do, yes — every claim on these slides is backed by a figure anyone can reproduce by running the notebooks, and ten automated tests confirm the underlying math is implemented correctly.",
    "It correctly tells apart real population differences from measurement gaps, and it shows the coverage promise breaking concretely, with real numbers, not just in theory.",
    "The finding holds up across two different clinical questions and two different scoring rules — it isn't a fluke of one particular choice.",
    "What's honestly still missing: reliability charts that go a layer deeper than coverage alone, coverage checks broken out by patient subgroup (not just by hospital), a flag for “was this value filled in or actually measured,” a first real attempt at fixing the problem rather than just detecting it, a genuinely new fifth hospital never touched during development, and mapping this dataset onto a real clinical data standard.",
], size=16.5, notes="\"Let's be honest with each other about where this actually stands. For what it "
         "set out to do — yes. Every claim on these slides is backed by a figure any of you can "
         "reproduce, right now, by running these same notebooks, and ten automated tests confirm the "
         "underlying math is implemented correctly. It correctly tells apart real population "
         "differences from measurement gaps, it showed the coverage promise breaking with real "
         "numbers, and it held up across two different clinical questions and two different scoring "
         "rules.\n\n"
         "\"But I want to be equally honest about what's still missing — deeper reliability checks, "
         "coverage broken out by patient subgroup and not just by hospital, a flag for whether a value "
         "was actually measured or just filled in, a real first attempt at fixing this rather than just "
         "detecting it, a genuinely new fifth hospital, and mapping this onto a real clinical data "
         "standard.\n\n"
         "\"Let me open the floor here: if you were adapting this pipeline for your own hospital's "
         "data, which of these missing pieces would you tackle first?\"")

# ---- References --------------------------------------------------------------
bullet_slide("References", "", [
    "A. N. Angelopoulos & S. Bates. “Conformal Prediction: A Gentle Introduction.” FTML 16(4), 2023. — the main source behind Section 8's math, in a genuinely readable form.",
    "Y. Romano, M. Sesia & E. Candès. “Classification with Valid and Adaptive Coverage.” NeurIPS 2020. — introduces the APS scoring rule used throughout this project.",
    "H. B. McMahan et al. “Communication-Efficient Learning of Deep Networks from Decentralized Data.” AISTATS 2017. — introduces FedAvg, the training method built from scratch in notebook 03.",
    "J. O. du Terrail et al. “FLamby: Datasets and Benchmarks for Cross-Silo Federated Learning.” 2022. — uses this same four-hospital population as a recognized benchmark.",
    "R. Detrano et al. “International Application of a New Probability Algorithm for the Diagnosis of Coronary Artery Disease.” Am. J. Cardiology, 1989. — the original clinical study behind this dataset.",
    "UCI Machine Learning Repository. “Heart Disease” dataset (ID 45). archive.ics.uci.edu/dataset/45/heart+disease.",
], size=14.5, notes="\"Quick pointer before we wrap up — if anyone wants to go deeper than we had time "
         "for today, start with Angelopoulos and Bates. It's written to be genuinely readable, no "
         "advanced statistics background assumed, and it's the source behind everything in Section "
         "8.\"")

# ---- Closing --------------------------------------------------------------
closing_slide(
    ["Questions & discussion", "fedconformal-clinical toolkit — Arnab Mukherjee",
     "Oklahoma State University — Center for Health Science, Tulsa"],
    notes="\"Let's close the loop on where we started. Three ideas, remember? Interoperability — "
          "getting a patient's data out of one hospital's system and into a form another system can "
          "read. You saw today that solving that alone gets you nothing about whether the data can be "
          "trusted once it arrives.\n\n"
          "\"Federated learning — several hospitals training one shared model together, without a "
          "single patient record ever leaving its own hospital. You built that yourselves this "
          "afternoon, in notebook 03, and saw it land within a point of the hypothetical pooled version "
          "on accuracy.\n\n"
          "\"And conformal prediction — the piece that actually checks whether the shared model can be "
          "trusted somewhere new. You built that too, in notebook 02, and then watched its promise "
          "silently break, from 90% down to 70.7%, the moment we deployed at a hospital it had never "
          "seen. That's the whole workshop: interoperability gets you the data, federated learning gets "
          "you the model, and conformal prediction is what tells you whether to trust it.\n\n"
          "\"If there's anything you didn't get to finish live, the exact reproduction command is on "
          "this slide, and the whole repository is yours to keep. Now — let's open it up. What "
          "questions do you have? [If the room is quiet, the 'what's honestly still missing' slide from "
          "the end of Section 9 is a good one to put back up and ask people to react to.] Thank you, "
          "everyone.\"",
)

# ---- Save -------------------------------------------------------------------
prs.save(OUT)
print(f"\n{_page_counter['n']} slides written to {os.path.abspath(OUT)}")
